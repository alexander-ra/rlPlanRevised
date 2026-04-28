"""Build the May 2026 Ruse University presentation deck.

Generates an editable .pptx from `presentation_script.md` using python-pptx.
The deck follows the script's design philosophy: clean scientific layout,
light background, dark charcoal text, restrained accents, custom shape-based
diagrams instead of stock photos.
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUTPUT = HERE / "adaptive_strategy_ruse_may.pptx"
LOGO = ASSETS / "ru-logo-125x140.png"


# -- color palette ---------------------------------------------------------

BG = RGBColor(0xFA, 0xFA, 0xFB)       # near-white background
INK = RGBColor(0x2C, 0x33, 0x3D)      # dark charcoal text
INK_SOFT = RGBColor(0x55, 0x60, 0x6D)  # secondary text
RULE = RGBColor(0xC9, 0xCE, 0xD6)     # divider line
MUTED = RGBColor(0x8A, 0x93, 0x9F)    # footer / faint marks

BLUE = RGBColor(0x1F, 0x4E, 0x79)     # adaptation / inference
BLUE_SOFT = RGBColor(0xDC, 0xE6, 0xF1)
RED = RGBColor(0xC0, 0x39, 0x2B)      # risk / exploitation
RED_SOFT = RGBColor(0xF7, 0xDD, 0xD9)
TEAL = RGBColor(0x16, 0x8B, 0x7C)     # evaluation / safety
TEAL_SOFT = RGBColor(0xD4, 0xEC, 0xE7)
AMBER = RGBColor(0xC7, 0x8A, 0x1F)    # accent for opponent types

FONT = "Calibri"


# -- low-level helpers -----------------------------------------------------


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor, width_pt: float = 1.0):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def no_line(shape):
    line = shape.line
    line.fill.background()


def add_text(
    shape,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    font: str = FONT,
):
    tf = shape.text_frame
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tf


def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    shape_type=MSO_SHAPE.RECTANGLE,
):
    shp = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill is None:
        shp.fill.background()
    else:
        set_fill(shp, fill)
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_width)
    shp.shadow.inherit = False
    return shp


def add_text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    **text_kwargs,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    add_text(box, text, **text_kwargs)
    return box


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = INK,
    width_pt: float = 1.0,
):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    return connector


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = INK,
    width_pt: float = 1.5,
):
    """Straight connector with an arrowhead at the (x2,y2) end."""
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(
        ln,
        qn("a:tailEnd"),
        {"type": "triangle", "w": "med", "len": "med"},
    )  # noqa: F841
    return conn


# -- slide-level scaffolding -----------------------------------------------


SLIDE_W = 13.333
SLIDE_H = 7.5


def setup_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # truly blank layout
    slide = prs.slides.add_slide(layout)
    bg = add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    bg.shadow.inherit = False
    return slide


def add_footer(slide, slide_number: int):
    """Thin divider, small logo + label on the left, slide number on the right."""
    add_line(slide, 0.5, 7.05, SLIDE_W - 0.5, 7.05, color=RULE, width_pt=0.75)
    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO), Inches(0.5), Inches(7.12), height=Inches(0.28)
        )
    add_text_box(
        slide,
        0.85,
        7.13,
        6.0,
        0.28,
        "Ruse University Academic Session  ·  May 2026",
        size=10,
        color=MUTED,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        SLIDE_W - 1.0,
        7.13,
        0.5,
        0.28,
        f"{slide_number:02d}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_slide_title(slide, title: str, *, accent: RGBColor = INK):
    add_text_box(
        slide,
        0.7,
        0.55,
        SLIDE_W - 1.4,
        0.7,
        title,
        size=30,
        bold=True,
        color=INK,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # short colored accent rule under the title
    bar = add_box(slide, 0.7, 1.28, 0.6, 0.06, fill=accent)
    bar.shadow.inherit = False


# -- Slide 1: Title --------------------------------------------------------


def slide_title(prs):
    slide = blank_slide(prs)

    # subtle right-side network pattern of nodes
    nodes = [
        (10.0, 2.0), (10.9, 2.7), (11.7, 2.1), (10.5, 3.5), (11.2, 3.9),
        (10.0, 4.6), (10.9, 5.0), (11.7, 4.5), (12.2, 3.2),
    ]
    # connections (faint)
    edges = [
        (0, 1), (1, 2), (1, 3), (3, 4), (4, 5), (4, 6), (6, 7), (7, 8),
        (2, 8), (3, 6), (1, 4), (5, 6),
    ]
    for a, b in edges:
        add_line(slide, nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1],
                 color=BLUE_SOFT, width_pt=1.0)
    for x, y in nodes:
        dot = add_box(slide, x - 0.09, y - 0.09, 0.18, 0.18,
                      fill=BLUE, shape_type=MSO_SHAPE.OVAL)
        dot.shadow.inherit = False

    # logo
    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO), Inches(0.7), Inches(0.7), height=Inches(1.2)
        )

    # main title
    add_text_box(
        slide, 0.7, 2.4, 8.6, 1.6,
        "Adaptive Strategy Learning\nin Multi-Agent\nImperfect-Information Games",
        size=36, bold=True, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    add_text_box(
        slide, 0.7, 4.55, 8.6, 0.5,
        "From Equilibrium Computation to Safe Opponent Exploitation",
        size=20, color=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # italicise subtitle
    box = slide.shapes[-1]
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.italic = True

    # accent rule
    bar = add_box(slide, 0.7, 5.25, 0.7, 0.07, fill=BLUE)
    bar.shadow.inherit = False

    add_text_box(
        slide, 0.7, 5.55, 8.6, 0.45,
        "Alexander Andreev",
        size=18, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide, 0.7, 6.0, 8.6, 0.4,
        "PhD session report  ·  May 2026  ·  Ruse University “Angel Kanchev”",
        size=14, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE,
    )


# -- Slide 2: Why this matters now -----------------------------------------


def slide_why_now(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Why this matters now", accent=BLUE)

    # bullets on the left
    bullets = [
        ("Strong single agents", "well understood in constrained benchmarks"),
        ("Agents interacting with other agents", "open problem"),
        ("Hidden information  ·  adversarial intent  ·  real-time decisions", None),
    ]
    y = 1.9
    for primary, secondary in bullets:
        # square bullet
        sq = add_box(slide, 0.85, y + 0.18, 0.16, 0.16, fill=BLUE)
        sq.shadow.inherit = False
        add_text_box(slide, 1.15, y, 6.0, 0.5, primary, size=18, bold=True, color=INK)
        if secondary:
            add_text_box(slide, 1.15, y + 0.5, 6.0, 0.4, secondary,
                         size=14, color=INK_SOFT)
        y += 1.05

    # right-side diagram: isolated agent vs. interacting agents
    # left panel header
    add_text_box(slide, 7.5, 1.7, 2.4, 0.4, "Single agent",
                 size=12, bold=True, color=INK_SOFT, align=PP_ALIGN.CENTER)
    # isolated blue node
    iso = add_box(slide, 8.45, 2.25, 0.5, 0.5, fill=BLUE,
                  shape_type=MSO_SHAPE.OVAL)
    iso.shadow.inherit = False
    add_text(iso, "A", size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)

    # divider
    add_line(slide, 10.4, 1.85, 10.4, 6.55, color=RULE, width_pt=0.75)

    # right panel header
    add_text_box(slide, 10.5, 1.7, 2.6, 0.4, "Multi-agent",
                 size=12, bold=True, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # multiple interacting agents
    nodes = [
        (10.95, 2.55, BLUE), (12.35, 2.55, RED),
        (10.65, 4.15, RED), (12.65, 4.15, BLUE),
        (11.65, 5.5, RED),
    ]
    labels = ["A", "?", "?", "?", "?"]
    placed = []
    for (x, y, c), lbl in zip(nodes, labels):
        n = add_box(slide, x - 0.25, y - 0.25, 0.5, 0.5, fill=c,
                    shape_type=MSO_SHAPE.OVAL)
        n.shadow.inherit = False
        add_text(n, lbl, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
        placed.append((x, y))
    # connections among them (dashed-style by using soft color)
    edges = [(0, 1), (0, 2), (1, 3), (2, 4), (3, 4), (1, 4), (0, 3)]
    for a, b in edges:
        add_line(slide, placed[a][0], placed[a][1],
                 placed[b][0], placed[b][1], color=RED, width_pt=1.25)

    add_text_box(slide, 10.45, 6.15, 2.7, 0.4,
                 "hidden  ·  possibly adversarial",
                 size=11, color=RED, align=PP_ALIGN.CENTER)

    add_footer(slide, 2)


# -- Slide 3: Use cases ----------------------------------------------------


def _icon_market(slide, cx, cy):
    """Stylised line chart inside a tile, anchored around (cx, cy)."""
    pts = [(-0.55, 0.25), (-0.30, 0.10), (-0.05, 0.20),
           (0.20, -0.10), (0.45, -0.05), (0.60, -0.25)]
    for i in range(len(pts) - 1):
        add_line(slide, cx + pts[i][0], cy + pts[i][1],
                 cx + pts[i + 1][0], cy + pts[i + 1][1],
                 color=BLUE, width_pt=2.0)
    # axes
    add_line(slide, cx - 0.6, cy + 0.35, cx - 0.6, cy - 0.35, color=INK_SOFT)
    add_line(slide, cx - 0.6, cy + 0.35, cx + 0.65, cy + 0.35, color=INK_SOFT)


def _icon_shield(slide, cx, cy):
    # stylised padlock: shackle (block arc, hollow) + body
    shackle = add_box(
        slide, cx - 0.22, cy - 0.45, 0.44, 0.55,
        fill=None, line=BLUE, line_width=2.5,
        shape_type=MSO_SHAPE.BLOCK_ARC,
    )
    shackle.shadow.inherit = False
    body = add_box(
        slide, cx - 0.30, cy - 0.10, 0.60, 0.55,
        fill=BLUE, line=BLUE, line_width=1.5,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    body.shadow.inherit = False
    # keyhole dot
    keyhole = add_box(
        slide, cx - 0.05, cy + 0.08, 0.10, 0.10,
        fill=BG, shape_type=MSO_SHAPE.OVAL,
    )
    keyhole.shadow.inherit = False


def _icon_network(slide, cx, cy):
    # central node + three orbiting nodes
    pts = [(-0.45, -0.25), (0.45, -0.2), (0.0, 0.4)]
    for x, y in pts:
        add_line(slide, cx, cy, cx + x, cy + y, color=BLUE, width_pt=1.5)
    for x, y in pts:
        n = add_box(slide, cx + x - 0.09, cy + y - 0.09, 0.18, 0.18,
                    fill=BLUE, shape_type=MSO_SHAPE.OVAL)
        n.shadow.inherit = False
    center = add_box(slide, cx - 0.13, cy - 0.13, 0.26, 0.26,
                     fill=RED, shape_type=MSO_SHAPE.OVAL)
    center.shadow.inherit = False


def _icon_patrol(slide, cx, cy):
    # grid 3x3 with two dots
    cell = 0.18
    origin_x = cx - 1.5 * cell
    origin_y = cy - 1.5 * cell
    for i in range(4):
        add_line(slide, origin_x + i * cell, origin_y,
                 origin_x + i * cell, origin_y + 3 * cell,
                 color=INK_SOFT, width_pt=0.75)
        add_line(slide, origin_x, origin_y + i * cell,
                 origin_x + 3 * cell, origin_y + i * cell,
                 color=INK_SOFT, width_pt=0.75)
    d1 = add_box(slide, origin_x + 0.05, origin_y + 0.05, 0.12, 0.12,
                 fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    d1.shadow.inherit = False
    d2 = add_box(slide, origin_x + 2 * cell + 0.05, origin_y + cell + 0.05,
                 0.12, 0.12, fill=RED, shape_type=MSO_SHAPE.OVAL)
    d2.shadow.inherit = False


def _icon_cards(slide, cx, cy):
    # two overlapping mini cards
    c1 = add_box(slide, cx - 0.35, cy - 0.30, 0.42, 0.55,
                 fill=BG, line=BLUE, line_width=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    c1.shadow.inherit = False
    add_text(c1, "♠", size=18, color=BLUE, align=PP_ALIGN.CENTER)
    c2 = add_box(slide, cx - 0.05, cy - 0.20, 0.42, 0.55,
                 fill=BG, line=RED, line_width=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    c2.shadow.inherit = False
    add_text(c2, "♥", size=18, color=RED, align=PP_ALIGN.CENTER)


def slide_use_cases(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Where this shows up", accent=BLUE)

    tiles = [
        ("Financial markets", "hidden positions  ·  microsecond reactions",
         _icon_market),
        ("Cybersecurity", "adaptive attackers  ·  hidden intent",
         _icon_shield),
        ("Social platforms", "bot networks  ·  coordinated disinformation",
         _icon_network),
        ("Security  (deployed today)", "patrol & screening randomisation",
         _icon_patrol),
        ("Gaming platforms", "collusion  ·  fraud detection",
         _icon_cards),
    ]
    # 5 tiles in two rows: 3 + 2 (centered)
    tile_w, tile_h = 3.9, 2.45
    gap = 0.25
    row1_y = 1.7
    row2_y = row1_y + tile_h + 0.15

    positions = [
        (0.5, row1_y),
        (0.5 + tile_w + gap, row1_y),
        (0.5 + 2 * (tile_w + gap), row1_y),
        (0.5 + 0.5 * (tile_w + gap), row2_y),
        (0.5 + 1.5 * (tile_w + gap), row2_y),
    ]

    for (title_txt, sub_txt, icon_fn), (left, top) in zip(tiles, positions):
        # tile background
        tile = add_box(slide, left, top, tile_w, tile_h,
                       fill=BG, line=RULE, line_width=0.75,
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        tile.shadow.inherit = False
        # icon area
        icon_cx = left + 0.85
        icon_cy = top + 1.05
        icon_fn(slide, icon_cx, icon_cy)
        # text area on the right of the tile
        add_text_box(slide, left + 1.55, top + 0.35, tile_w - 1.7, 0.55,
                     title_txt, size=15, bold=True, color=INK,
                     anchor=MSO_ANCHOR.TOP)
        add_text_box(slide, left + 1.55, top + 0.95, tile_w - 1.7, 1.4,
                     sub_txt, size=11, color=INK_SOFT,
                     anchor=MSO_ANCHOR.TOP)

    add_footer(slide, 3)


# -- Slide 4: Why study games? ---------------------------------------------


def slide_testbeds(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Why study games?", accent=BLUE)

    add_text_box(slide, 0.7, 1.4, SLIDE_W - 1.4, 0.4,
                 "Controlled environments where strategic behaviour is testable.",
                 size=14, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)

    panels = [
        ("Poker  /  Belot-like",
         "hidden cards  ·  inference",
         "cards"),
        ("Auction-style",
         "hidden valuations  ·  bidding",
         "gavel"),
        ("Pursuit-evasion",
         "partial observability  ·  adversarial search",
         "grid"),
        ("Coalition games",
         "alliances  ·  betrayal  ·  N-player incentives",
         "graph"),
    ]
    panel_w = 2.85
    panel_h = 4.6
    gap = 0.2
    total = 4 * panel_w + 3 * gap
    start_x = (SLIDE_W - total) / 2
    top = 2.05

    for i, (title_txt, sub_txt, kind) in enumerate(panels):
        left = start_x + i * (panel_w + gap)
        panel = add_box(slide, left, top, panel_w, panel_h,
                        fill=BG, line=RULE, line_width=0.75,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        panel.shadow.inherit = False

        cx = left + panel_w / 2
        cy = top + 1.4
        if kind == "cards":
            _icon_cards(slide, cx, cy)
            # corner glyph - small playful card pip in upper-left of panel
            add_text_box(slide, left + 0.15, top + 0.1, 0.4, 0.4, "♣",
                         size=12, color=MUTED)
        elif kind == "gavel":
            # auction: vertical gavel + base
            add_line(slide, cx - 0.35, cy + 0.25, cx + 0.35, cy + 0.25,
                     color=BLUE, width_pt=3.5)
            add_line(slide, cx + 0.05, cy + 0.25, cx + 0.45, cy - 0.30,
                     color=BLUE, width_pt=3.0)
            add_line(slide, cx - 0.45, cy + 0.45, cx + 0.45, cy + 0.45,
                     color=INK_SOFT, width_pt=2.0)
        elif kind == "grid":
            # grid + pursuer (red) + evader (blue)
            cell = 0.20
            ox = cx - 2 * cell
            oy = cy - 2 * cell
            for k in range(5):
                add_line(slide, ox + k * cell, oy, ox + k * cell, oy + 4 * cell,
                         color=INK_SOFT, width_pt=0.75)
                add_line(slide, ox, oy + k * cell, ox + 4 * cell, oy + k * cell,
                         color=INK_SOFT, width_pt=0.75)
            p = add_box(slide, ox + 0.05, oy + 0.05, 0.13, 0.13,
                        fill=RED, shape_type=MSO_SHAPE.OVAL)
            p.shadow.inherit = False
            e = add_box(slide, ox + 3 * cell + 0.05, oy + 3 * cell + 0.05,
                        0.13, 0.13, fill=BLUE, shape_type=MSO_SHAPE.OVAL)
            e.shadow.inherit = False
        elif kind == "graph":
            # 4 nodes with alliance lines (solid blue) and betrayal (red dashed)
            pts = [(cx - 0.45, cy - 0.30), (cx + 0.45, cy - 0.30),
                   (cx + 0.45, cy + 0.30), (cx - 0.45, cy + 0.30)]
            edges = [(0, 1, BLUE), (1, 2, BLUE), (0, 3, BLUE), (0, 2, RED)]
            for a, b, col in edges:
                add_line(slide, pts[a][0], pts[a][1],
                         pts[b][0], pts[b][1], color=col, width_pt=1.5)
            for x, y in pts:
                n = add_box(slide, x - 0.10, y - 0.10, 0.20, 0.20,
                            fill=BLUE, shape_type=MSO_SHAPE.OVAL)
                n.shadow.inherit = False

        # panel text
        add_text_box(slide, left + 0.2, top + 2.55, panel_w - 0.4, 0.5,
                     title_txt, size=14, bold=True, color=INK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, left + 0.2, top + 3.1, panel_w - 0.4, 1.3,
                     sub_txt, size=11, color=INK_SOFT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    add_text_box(slide, 0.7, 6.7, SLIDE_W - 1.4, 0.3,
                 "Validation environments — not the final application domains.",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 4)


# -- Slide 5: State of the art --------------------------------------------


def slide_state_of_the_art(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "State of the art", accent=BLUE)

    add_text_box(slide, 0.7, 1.35, SLIDE_W - 1.4, 0.4,
                 "“Superhuman” in progressively larger games.",
                 size=14, color=INK_SOFT)

    # horizontal timeline
    y_axis = 4.0
    x_left = 1.1
    x_right = SLIDE_W - 1.1
    add_line(slide, x_left, y_axis, x_right, y_axis, color=INK_SOFT, width_pt=2.0)
    # arrow at right end
    add_line(slide, x_right - 0.2, y_axis - 0.12, x_right, y_axis,
             color=INK_SOFT, width_pt=2.0)
    add_line(slide, x_right - 0.2, y_axis + 0.12, x_right, y_axis,
             color=INK_SOFT, width_pt=2.0)

    milestones = [
        ("2007", "CFR", "regret minimisation\nin imperfect-info games"),
        ("2017", "Libratus", "2-player poker\nsuperhuman"),
        ("2019", "Pluribus", "6-player poker\nsuperhuman"),
        ("2020", "ReBeL", "search + RL\nin imperfect-info"),
        ("2022", "DeepNash  /  CICERO", "Stratego  ·  Diplomacy\n(7 players, coalitions)"),
    ]
    n = len(milestones)
    span = x_right - x_left - 0.6
    for i, (year, name, sub) in enumerate(milestones):
        x = x_left + 0.3 + i * span / (n - 1)
        # marker dot
        dot = add_box(slide, x - 0.13, y_axis - 0.13, 0.26, 0.26,
                      fill=BLUE, shape_type=MSO_SHAPE.OVAL)
        dot.shadow.inherit = False
        # alternate above/below
        above = (i % 2 == 0)
        if above:
            add_text_box(slide, x - 1.1, y_axis - 1.7, 2.2, 0.4, year,
                         size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
            add_text_box(slide, x - 1.1, y_axis - 1.3, 2.2, 0.5, name,
                         size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
            add_text_box(slide, x - 1.3, y_axis - 0.78, 2.6, 0.6, sub,
                         size=10, color=INK_SOFT, align=PP_ALIGN.CENTER)
            add_line(slide, x, y_axis - 0.1, x, y_axis - 0.45,
                     color=BLUE, width_pt=1.5)
        else:
            add_line(slide, x, y_axis + 0.1, x, y_axis + 0.45,
                     color=BLUE, width_pt=1.5)
            add_text_box(slide, x - 1.1, y_axis + 0.5, 2.2, 0.4, year,
                         size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
            add_text_box(slide, x - 1.1, y_axis + 0.9, 2.2, 0.5, name,
                         size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
            add_text_box(slide, x - 1.3, y_axis + 1.4, 2.6, 0.6, sub,
                         size=10, color=INK_SOFT, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)


# -- Slide 6: The limitation -----------------------------------------------


def slide_limitation(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "The limitation", accent=RED)

    # central heavy statement
    add_text_box(slide, 0.7, 1.55, SLIDE_W - 1.4, 0.7,
                 "Fixed strategies. No adaptation.",
                 size=30, bold=True, color=RED, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, 0.7, 2.25, SLIDE_W - 1.4, 0.4,
                 "The same policy against every opponent.",
                 size=16, color=INK_SOFT, align=PP_ALIGN.LEFT)

    # central policy box
    cx, cy = 3.1, 4.85
    pol = add_box(slide, cx - 1.2, cy - 0.6, 2.4, 1.2,
                  fill=BG, line=INK, line_width=2.0,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    pol.shadow.inherit = False
    add_text(pol, "Equilibrium\npolicy",
             size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # three opponent boxes on the right with identical arrows
    opps = [
        ("Cautious beginner", BLUE),
        ("Aggressive bluffer", RED),
        ("Coordinated group", AMBER),
    ]
    box_left = 8.4
    box_w = 3.4
    box_h = 1.0
    gap = 0.4
    total_h = 3 * box_h + 2 * gap
    start_y = cy - total_h / 2
    for i, (label, col) in enumerate(opps):
        top = start_y + i * (box_h + gap)
        b = add_box(slide, box_left, top, box_w, box_h,
                    fill=BG, line=col, line_width=2.0,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        b.shadow.inherit = False
        add_text(b, label, size=14, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        # identical arrow from policy to this opponent
        add_arrow(slide,
                  cx + 1.2, cy,
                  box_left, top + box_h / 2,
                  color=INK_SOFT, width_pt=1.75)

    # caption
    add_text_box(slide, 0.7, 6.55, SLIDE_W - 1.4, 0.4,
                 "Safe — but blind. Systematically misses exploitable weaknesses.",
                 size=12, color=MUTED, align=PP_ALIGN.LEFT)

    add_footer(slide, 6)


# -- Slide 7: Three open problems -----------------------------------------


def slide_three_problems(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Three open problems", accent=INK)

    # left: numbered list
    items = [
        ("1", "Adaptation", "infer opponents in real time", BLUE),
        ("2", "Safety", "exploit without exposing yourself", RED),
        ("3", "Evaluation", "measure if an agent adapts well", TEAL),
    ]
    y = 2.0
    for num, name, sub, col in items:
        circ = add_box(slide, 0.85, y, 0.7, 0.7, fill=col,
                       shape_type=MSO_SHAPE.OVAL)
        circ.shadow.inherit = False
        add_text(circ, num, size=22, bold=True, color=BG,
                 align=PP_ALIGN.CENTER)
        add_text_box(slide, 1.7, y - 0.05, 5.5, 0.5, name,
                     size=22, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, 1.7, y + 0.4, 5.5, 0.45, sub,
                     size=13, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.4

    # right: triangle / loop diagram
    cx, cy, r = 10.6, 4.4, 1.55
    # vertices of triangle pointing up
    import math
    verts = []
    for k in range(3):
        ang = math.radians(-90 + k * 120)
        verts.append((cx + r * math.cos(ang), cy + r * math.sin(ang)))
    cols = [BLUE, RED, TEAL]
    labels = ["Adaptation", "Safety", "Evaluation"]

    # connecting curved-ish arrows (use straight arrows)
    for i in range(3):
        a = verts[i]
        b = verts[(i + 1) % 3]
        # shorten slightly so arrow doesn't enter circle
        dx, dy = b[0] - a[0], b[1] - a[1]
        L = (dx ** 2 + dy ** 2) ** 0.5
        ux, uy = dx / L, dy / L
        sa = (a[0] + 0.45 * ux, a[1] + 0.45 * uy)
        sb = (b[0] - 0.45 * ux, b[1] - 0.45 * uy)
        add_arrow(slide, sa[0], sa[1], sb[0], sb[1], color=INK_SOFT, width_pt=1.5)

    for (vx, vy), col, lbl in zip(verts, cols, labels):
        n = add_box(slide, vx - 0.55, vy - 0.55, 1.1, 1.1, fill=col,
                    shape_type=MSO_SHAPE.OVAL)
        n.shadow.inherit = False
        add_text(n, lbl, size=11, bold=True, color=BG,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, 8.7, 6.45, 4.0, 0.35,
                 "inference  →  action  →  measurement",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 7)


# -- Slide 8: Contribution 1 — Behavioral Adaptation -----------------------


def slide_contrib_1(prs):
    slide = blank_slide(prs)
    # eyebrow label + title
    add_text_box(slide, 0.7, 0.45, 6.0, 0.35,
                 "CONTRIBUTION  1",
                 size=11, bold=True, color=BLUE)
    add_slide_title(slide, "Behavioral Adaptation Framework", accent=BLUE)

    # left bullets
    bullets = [
        ("Infer opponent type", "not just hidden state"),
        ("Adapt only when evidence is strong enough", "fallback to equilibrium"),
        ("Detect anomalies", "bots  ·  collusion  ·  adversarial users"),
    ]
    y = 2.0
    for primary, secondary in bullets:
        sq = add_box(slide, 0.85, y + 0.18, 0.16, 0.16, fill=BLUE)
        sq.shadow.inherit = False
        add_text_box(slide, 1.15, y, 5.5, 0.5, primary,
                     size=17, bold=True, color=INK)
        add_text_box(slide, 1.15, y + 0.5, 5.5, 0.4, secondary,
                     size=12, color=INK_SOFT)
        y += 1.25

    # right: belief-state diagram
    # 1) "Observed actions" source on the left of the diagram column
    src_x, src_y = 7.4, 4.0
    src = add_box(slide, src_x, src_y - 0.4, 1.7, 0.8,
                  fill=BG, line=INK, line_width=1.5,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    src.shadow.inherit = False
    add_text(src, "Observed\nactions", size=12, bold=True, color=INK,
             align=PP_ALIGN.CENTER)

    # 2) two belief boxes
    b1 = add_box(slide, 9.6, 2.4, 2.2, 0.85,
                 fill=BLUE_SOFT, line=BLUE, line_width=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    b1.shadow.inherit = False
    add_text(b1, "Belief: hidden state",
             size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    b2 = add_box(slide, 9.6, 4.7, 2.2, 0.85,
                 fill=BLUE_SOFT, line=BLUE, line_width=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    b2.shadow.inherit = False
    add_text(b2, "Belief: behavioural type",
             size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # arrows from source to beliefs
    add_arrow(slide, src_x + 1.7, src_y - 0.15, 9.6, 2.85,
              color=INK_SOFT, width_pt=1.5)
    add_arrow(slide, src_x + 1.7, src_y + 0.15, 9.6, 5.1,
              color=INK_SOFT, width_pt=1.5)

    # 3) decision-making
    dec = add_box(slide, 12.0, 3.6, 1.0, 1.0,
                  fill=BLUE, line=BLUE, line_width=1.5,
                  shape_type=MSO_SHAPE.OVAL)
    dec.shadow.inherit = False
    add_text(dec, "Decide", size=11, bold=True, color=BG,
             align=PP_ALIGN.CENTER)

    # arrows from beliefs to decision
    add_arrow(slide, 11.8, 2.85, 12.4, 3.6, color=INK_SOFT, width_pt=1.5)
    add_arrow(slide, 11.8, 5.1, 12.4, 4.6, color=INK_SOFT, width_pt=1.5)

    # 4) fallback arrow when evidence is weak
    fb = add_box(slide, 9.6, 6.0, 2.2, 0.55,
                 fill=BG, line=MUTED, line_width=1.0,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    fb.shadow.inherit = False
    add_text(fb, "Fallback: equilibrium play",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # weak-evidence arrow from source to fallback
    conn = slide.shapes.add_connector(
        2,  # bent connector
        Inches(src_x + 0.85), Inches(src_y + 0.4),
        Inches(10.7), Inches(6.0),
    )
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(1.0)
    add_text_box(slide, 7.5, 5.55, 2.0, 0.3,
                 "if evidence weak", size=10,
                 color=MUTED)

    # tiny bot icon near anomaly detection (top-right corner of diagram)
    bot_x, bot_y = 12.4, 1.65
    head = add_box(slide, bot_x - 0.18, bot_y - 0.05, 0.36, 0.32,
                   fill=BG, line=BLUE, line_width=1.5,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    head.shadow.inherit = False
    eye1 = add_box(slide, bot_x - 0.10, bot_y + 0.05, 0.07, 0.07,
                   fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    eye1.shadow.inherit = False
    eye2 = add_box(slide, bot_x + 0.03, bot_y + 0.05, 0.07, 0.07,
                   fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    eye2.shadow.inherit = False
    # antenna
    add_line(slide, bot_x, bot_y - 0.05, bot_x, bot_y - 0.20,
             color=BLUE, width_pt=1.5)
    ant = add_box(slide, bot_x - 0.04, bot_y - 0.27, 0.08, 0.08,
                  fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    ant.shadow.inherit = False
    add_text_box(slide, 11.7, 2.0, 1.7, 0.3, "anomaly",
                 size=9, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 8)


# -- Slide 9: Contribution 2 — Safe Exploitation ---------------------------


def slide_contrib_2(prs):
    slide = blank_slide(prs)
    add_text_box(slide, 0.7, 0.45, 6.0, 0.35,
                 "CONTRIBUTION  2",
                 size=11, bold=True, color=RED)
    add_slide_title(slide, "Multi-Agent Safe Exploitation", accent=RED)

    # left bullets
    bullets = [
        ("Exploit detected weaknesses", "act on the inferred type"),
        ("Stay close to a reference policy", "πKL  regularisation"),
        ("Useful safety heuristics", "beyond two-player guarantees"),
    ]
    y = 2.0
    for primary, secondary in bullets:
        sq = add_box(slide, 0.85, y + 0.18, 0.16, 0.16, fill=RED)
        sq.shadow.inherit = False
        add_text_box(slide, 1.15, y, 5.5, 0.5, primary,
                     size=17, bold=True, color=INK)
        add_text_box(slide, 1.15, y + 0.5, 5.5, 0.4, secondary,
                     size=12, color=INK_SOFT)
        y += 1.25

    # right: dial / slider from "exploit" to "reference policy"
    dial_left = 7.6
    dial_right = 12.9
    dial_y = 3.4

    # endpoint labels
    add_text_box(slide, dial_left - 0.4, dial_y - 1.2, 2.2, 0.4,
                 "Exploit", size=14, bold=True, color=RED,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, dial_right - 1.8, dial_y - 1.2, 2.2, 0.4,
                 "Reference policy", size=14, bold=True, color=BLUE,
                 align=PP_ALIGN.CENTER)

    # track
    track = add_box(slide, dial_left, dial_y - 0.10, dial_right - dial_left,
                    0.20, fill=BLUE_SOFT, line=None,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    track.shadow.inherit = False

    # gradient-ish: a colored band from red side → middle
    colored = add_box(slide, dial_left, dial_y - 0.10,
                      (dial_right - dial_left) * 0.45, 0.20,
                      fill=RED_SOFT, line=None,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    colored.shadow.inherit = False

    # tick marks
    for k in range(5):
        x = dial_left + k * (dial_right - dial_left) / 4
        add_line(slide, x, dial_y - 0.25, x, dial_y - 0.10,
                 color=INK_SOFT, width_pt=1.0)

    # knob a third of the way from "exploit"
    knob_x = dial_left + 0.35 * (dial_right - dial_left)
    knob = add_box(slide, knob_x - 0.18, dial_y - 0.32, 0.36, 0.65,
                   fill=INK, line=None,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    knob.shadow.inherit = False

    # "constraint band" annotation
    add_text_box(slide, dial_left, dial_y + 0.55, dial_right - dial_left, 0.4,
                 "constrained by  πKL  band around reference",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # bottom note
    note = add_box(slide, 7.4, 5.0, 5.7, 1.5,
                   fill=BG, line=RULE, line_width=0.75,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    note.shadow.inherit = False
    add_text(note,
             "Reference policy:\n"
             "population average  ·  human-behaviour prior  ·  approximate equilibrium",
             size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # 2015 safety theorem callout (small note on left bottom)
    callout = add_box(slide, 0.85, 5.6, 6.2, 1.0,
                      fill=RED_SOFT, line=RED, line_width=1.0,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    callout.shadow.inherit = False
    add_text(callout,
             "Two-player safety guarantees provably fail at  N ≥ 3.",
             size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

    add_footer(slide, 9)


# -- Slide 10: Contribution 3 — Evaluation --------------------------------


def slide_contrib_3(prs):
    slide = blank_slide(prs)
    add_text_box(slide, 0.7, 0.45, 6.0, 0.35,
                 "CONTRIBUTION  3",
                 size=11, bold=True, color=TEAL)
    add_slide_title(slide, "Evaluation Methodology", accent=TEAL)

    # left bullets
    bullets = [
        ("Safety", "worst-case vulnerability"),
        ("Population ranking", "non-transitive dynamics"),
        ("Statistical reliability", "variance reduction"),
    ]
    y = 2.0
    for primary, secondary in bullets:
        sq = add_box(slide, 0.85, y + 0.18, 0.16, 0.16, fill=TEAL)
        sq.shadow.inherit = False
        add_text_box(slide, 1.15, y, 5.5, 0.5, primary,
                     size=17, bold=True, color=INK)
        add_text_box(slide, 1.15, y + 0.5, 5.5, 0.4, secondary,
                     size=12, color=INK_SOFT)
        y += 1.25

    # right: three-layer evaluation stack
    stack_left = 8.0
    stack_w = 4.6
    layer_h = 1.0
    base_y = 5.7
    layers = [
        ("Statistical reliability", "AIVAT", TEAL),
        ("Population ranking",      "α-Rank", TEAL),
        ("Safety",                  "N-player exploitability", TEAL),
    ]
    for i, (lbl, method, col) in enumerate(layers):
        top = base_y - (i + 1) * layer_h
        # alternate alpha by tinting - use TEAL_SOFT for top, slightly stronger lower
        fill = TEAL_SOFT if i != 1 else BG
        layer = add_box(slide, stack_left, top, stack_w, layer_h - 0.08,
                        fill=fill, line=col, line_width=1.5,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        layer.shadow.inherit = False
        add_text_box(slide, stack_left + 0.2, top + 0.1,
                     stack_w - 0.4, 0.35, lbl,
                     size=14, bold=True, color=INK)
        add_text_box(slide, stack_left + 0.2, top + 0.45,
                     stack_w - 0.4, 0.35, method,
                     size=12, color=col)

    # rock-paper-scissors small glyphs next to "Population ranking" layer
    rps_y = base_y - 2 * layer_h + 0.45
    rps_x = stack_left + stack_w + 0.15
    add_text_box(slide, rps_x, rps_y - 0.15, 0.4, 0.4, "✊",
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, rps_x + 0.30, rps_y - 0.15, 0.4, 0.4, "✋",
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, rps_x + 0.60, rps_y - 0.15, 0.4, 0.4, "✌",
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_text_box(slide, stack_left, 5.85, stack_w, 0.35,
                 "validated across structurally different game types",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 10)


# -- Slide 11: Expected outcomes ------------------------------------------


def slide_outcomes(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Expected outcomes", accent=INK)

    col_top = 1.85
    col_h = 4.7
    gap = 0.4
    col_w = (SLIDE_W - 1.4 - gap) / 2
    left_x = 0.7
    right_x = left_x + col_w + gap

    # left column — realistic target (grounded)
    left = add_box(slide, left_x, col_top, col_w, col_h,
                   fill=BG, line=BLUE, line_width=2.0,
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    left.shadow.inherit = False
    add_text_box(slide, left_x + 0.3, col_top + 0.25, col_w - 0.6, 0.4,
                 "REALISTIC TARGET", size=11, bold=True, color=BLUE)
    add_text_box(slide, left_x + 0.3, col_top + 0.7, col_w - 0.6, 0.6,
                 "What the thesis will deliver",
                 size=18, bold=True, color=INK)

    items_left = [
        "Validated adaptation methods",
        "Characterised safety heuristics",
        "Working evaluation framework",
        "Failure modes documented",
    ]
    y = col_top + 1.6
    for txt in items_left:
        sq = add_box(slide, left_x + 0.4, y + 0.18, 0.16, 0.16, fill=BLUE)
        sq.shadow.inherit = False
        add_text_box(slide, left_x + 0.7, y, col_w - 0.9, 0.5, txt,
                     size=14, color=INK)
        y += 0.7

    # right column — stretch (lighter)
    right = add_box(slide, right_x, col_top, col_w, col_h,
                    fill=BG, line=RULE, line_width=1.0,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    right.shadow.inherit = False
    add_text_box(slide, right_x + 0.3, col_top + 0.25, col_w - 0.6, 0.4,
                 "STRETCH DIRECTION", size=11, bold=True, color=MUTED)
    add_text_box(slide, right_x + 0.3, col_top + 0.7, col_w - 0.6, 0.6,
                 "If the work goes deeper",
                 size=18, bold=True, color=INK_SOFT)

    items_right = [
        "Theoretical extensions of the safety theorem",
        "Novel combinations across contributions",
        "Surprising empirical results",
    ]
    y = col_top + 1.6
    for txt in items_right:
        sq = add_box(slide, right_x + 0.4, y + 0.22, 0.12, 0.12,
                     fill=MUTED, shape_type=MSO_SHAPE.OVAL)
        sq.shadow.inherit = False
        add_text_box(slide, right_x + 0.7, y, col_w - 0.9, 0.5, txt,
                     size=13, color=INK_SOFT)
        y += 0.7

    add_text_box(slide, 0.7, 6.55, SLIDE_W - 1.4, 0.4,
                 "Grounded outcomes already constitute a coherent thesis. No pivot required.",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 11)


# -- Slide 12: Close ------------------------------------------------------


def slide_close(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Today  →  Tomorrow", accent=INK)

    add_text_box(slide, 0.7, 1.45, SLIDE_W - 1.4, 0.5,
                 "The gap between today’s deployed systems and the next generation.",
                 size=18, color=INK_SOFT)

    # left side: today
    left_x, left_y, w, h = 1.0, 2.7, 4.6, 3.4
    today = add_box(slide, left_x, left_y, w, h,
                    fill=BG, line=RULE, line_width=1.0,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    today.shadow.inherit = False
    add_text_box(slide, left_x + 0.3, left_y + 0.3, w - 0.6, 0.4,
                 "TODAY", size=11, bold=True, color=MUTED)
    add_text_box(slide, left_x + 0.3, left_y + 0.7, w - 0.6, 0.5,
                 "Fixed-strategy systems",
                 size=20, bold=True, color=INK)
    items = [
        "Equilibrium solvers",
        "Identical play vs. every opponent",
        "Two-player safety guarantees",
    ]
    y = left_y + 1.5
    for txt in items:
        sq = add_box(slide, left_x + 0.4, y + 0.18, 0.14, 0.14,
                     fill=MUTED, shape_type=MSO_SHAPE.RECTANGLE)
        sq.shadow.inherit = False
        add_text_box(slide, left_x + 0.7, y, w - 0.9, 0.5, txt,
                     size=13, color=INK_SOFT)
        y += 0.55

    # arrow in the middle
    add_arrow(slide, 5.9, 4.4, 7.5, 4.4, color=INK, width_pt=2.5)
    add_text_box(slide, 6.0, 4.45, 1.5, 0.35, "thesis",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # right side: tomorrow
    right_x = 7.7
    tomorrow = add_box(slide, right_x, left_y, w, h,
                       fill=BG, line=BLUE, line_width=2.0,
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tomorrow.shadow.inherit = False
    add_text_box(slide, right_x + 0.3, left_y + 0.3, w - 0.6, 0.4,
                 "NEXT GENERATION", size=11, bold=True, color=BLUE)
    add_text_box(slide, right_x + 0.3, left_y + 0.7, w - 0.6, 0.5,
                 "Adaptive, safe, accountable agents",
                 size=20, bold=True, color=INK)
    items_r = [
        ("Adaptation", BLUE),
        ("Safety beyond two players", RED),
        ("Systematic evaluation", TEAL),
    ]
    y = left_y + 1.5
    for txt, col in items_r:
        sq = add_box(slide, right_x + 0.4, y + 0.18, 0.14, 0.14, fill=col)
        sq.shadow.inherit = False
        add_text_box(slide, right_x + 0.7, y, w - 0.9, 0.5, txt,
                     size=13, bold=True, color=INK)
        y += 0.55

    # closing thank-you line
    add_text_box(slide, 0.7, 6.45, SLIDE_W - 1.4, 0.5,
                 "Thank you.   Questions welcome.",
                 size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_footer(slide, 12)


# -- Build entry point ----------------------------------------------------


def build():
    prs = setup_presentation()
    slide_title(prs)
    slide_why_now(prs)
    slide_use_cases(prs)
    slide_testbeds(prs)
    slide_state_of_the_art(prs)
    slide_limitation(prs)
    slide_three_problems(prs)
    slide_contrib_1(prs)
    slide_contrib_2(prs)
    slide_contrib_3(prs)
    slide_outcomes(prs)
    slide_close(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
