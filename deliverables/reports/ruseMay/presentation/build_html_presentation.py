"""Build the May 2026 Ruse University deck from HTML slide designs.

Pipeline:
1. Write one standalone 1920x1080 HTML file per slide.
2. Render each HTML slide to a 3840x2160 PNG with Playwright/Chromium
   using device_scale_factor=2 for sharper projected text.
3. Place each PNG full-bleed into a 16:9 PowerPoint deck.

This intentionally makes the PPTX non-editable at the element level. The
editable source of truth is the generated HTML/CSS in ``html_slides/``.
"""

from __future__ import annotations

import base64
import html
import json
import mimetypes
import re
import shutil
import subprocess
import textwrap
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.util import Inches


HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
LOGO = ASSETS / "ru-logo-125x140.png"

HTML_DIR = HERE / "html_slides"
PNG_DIR = HERE / "slide_images"
OUTPUT = HERE / "adaptive_strategy_ruse_may_html.pptx"
SCRIPT_PATH = HERE.parent / "presentation_script.md"

WIDTH = 1920
HEIGHT = 1080


def logo_data_uri() -> str:
    if not LOGO.exists():
        return ""
    data = base64.b64encode(LOGO.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{data}"


LOGO_URI = logo_data_uri()


def asset_data_uri(relative_path: str) -> str:
    path = ASSETS / relative_path
    if not path.exists():
        return ""
    mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{data}"


def asset_img(relative_path: str, alt: str, class_name: str) -> str:
    uri = asset_data_uri(relative_path)
    if not uri:
        return ""
    return f'<img class="{esc(class_name)}" src="{uri}" alt="{esc(alt)}">'


def thumbnail(relative_path: str, alt: str) -> str:
    image = asset_img(relative_path, alt, "example-thumb")
    if image:
        return image
    return f'<div class="example-thumb placeholder-thumb">{esc(alt)}</div>'


def esc(value: str) -> str:
    return html.escape(value, quote=True)


TOTAL_SLIDES = 13


def footer(slide_no: int, title: str) -> str:
    return f"""
    <footer class="footer">
      <div class="footer-left">
        {'<img src="' + LOGO_URI + '" alt="University of Ruse logo">' if LOGO_URI else ''}
        <span>&ldquo;Angel Kanchev&rdquo; University of Ruse</span>
      </div>
      <div class="footer-right">
        <span class="footer-title">Alexander Andreev &middot; Prof. Dr. Tsvetomir Vasilev</span>
        <span class="slide-no">{slide_no:02d}/{TOTAL_SLIDES}</span>
      </div>
    </footer>
    """


COMMON_CSS = """\
/* === 1. Design tokens === */
:root {
  --bg: #f5f7fb;
  --panel: #ffffff;
  --ink: #17212c;
  --ink-soft: #526173;
  --muted: #8793a5;
  --rule: #d7dde7;
  --blue: #155f9f;
  --blue-soft: #dbeaf8;
  --red: #c74735;
  --red-soft: #fae3de;
  --teal: #148d80;
  --teal-soft: #d8eee9;
  --amber: #bb7d1c;
  --shadow: 0 26px 70px rgba(22, 34, 50, 0.10);
  --shadow-soft: 0 12px 32px rgba(22, 34, 50, 0.07);
  --hairline: rgba(21, 95, 159, 0.12);
}

/* === 2. Reset === */
* { box-sizing: border-box; }
html, body {
  margin: 0;
  width: 1920px;
  height: 1080px;
  overflow: hidden;
  background: var(--bg);
  color: var(--ink);
  font-family: Aptos, "Segoe UI", Arial, Helvetica, sans-serif;
  letter-spacing: 0;
}
body { -webkit-font-smoothing: antialiased; }

/* === 3. Slide canvas === */
.slide {
  position: relative;
  width: 1920px;
  height: 1080px;
  padding: 82px 104px 104px;
  overflow: hidden;
  background:
    linear-gradient(180deg, rgba(255,255,255,0.94), rgba(245,247,251,0.98)),
    linear-gradient(135deg, rgba(21,95,159,0.06), transparent 28%),
    linear-gradient(315deg, rgba(20,141,128,0.055), transparent 24%),
    var(--bg);
}
.slide::before {
  content: "";
  position: absolute;
  inset: 0;
  z-index: 0;
  pointer-events: none;
  background-image:
    linear-gradient(var(--hairline) 1px, transparent 1px),
    linear-gradient(90deg, var(--hairline) 1px, transparent 1px);
  background-size: 64px 64px;
  mask-image: linear-gradient(90deg, transparent, #000 14%, #000 86%, transparent);
  opacity: 0.38;
}
.slide::after {
  content: "";
  position: absolute;
  left: 104px;
  right: 104px;
  top: 30px;
  height: 4px;
  z-index: 0;
  border-radius: 999px;
  background: linear-gradient(90deg, var(--blue), var(--teal), transparent 72%);
  opacity: 0.55;
}
.slide > * {
  position: relative;
  z-index: 1;
}
.title-slide {
  background:
    linear-gradient(90deg, rgba(255,255,255,0.98), rgba(245,247,251,0.9) 58%, rgba(245,247,251,0.98)),
    linear-gradient(135deg, rgba(21,95,159,0.08), transparent 34%),
    linear-gradient(315deg, rgba(20,141,128,0.07), transparent 25%),
    var(--bg);
}
.title-slide::after {
  left: 104px;
  right: 104px;
  background: linear-gradient(90deg, var(--blue), var(--teal), var(--red), transparent 76%);
}

/* === 4. Header & title rule === */
.slide-header {
  position: relative;
  height: 128px;
}
h1 {
  margin: 0;
  font-size: 56px;
  line-height: 1.04;
  font-weight: 780;
  max-width: 1300px;
  letter-spacing: -0.01em;
}
.title-rule {
  width: 112px;
  height: 8px;
  margin-top: 24px;
  border-radius: 999px;
  background: var(--ink);
  box-shadow: 0 6px 16px rgba(21, 95, 159, 0.18);
}
.title-rule.blue  { background: var(--blue); }
.title-rule.red   { background: var(--red); }
.title-rule.teal  { background: var(--teal); }

/* === 5. Footer === */
.footer {
  position: absolute;
  left: 104px;
  right: 104px;
  bottom: 36px;
  height: 42px;
  border-top: 1.5px solid var(--rule);
  display: flex;
  align-items: flex-end;
  justify-content: space-between;
  color: var(--muted);
  font-size: 18px;
  background: linear-gradient(180deg, transparent, rgba(245,247,251,0.75) 42%);
}
.footer-left { display: flex; align-items: center; gap: 14px; }
.footer img  { width: 30px; height: auto; display: block; }
.footer-right {
  display: flex;
  align-items: flex-end;
  gap: 22px;
}
.footer-title {
  color: var(--ink-soft);
  font-weight: 700;
  max-width: 720px;
  overflow: hidden;
  text-overflow: ellipsis;
  white-space: nowrap;
}
.slide-no    { font-variant-numeric: tabular-nums; }

/* === 6. Typography === */
.eyebrow {
  margin-bottom: 16px;
  color: var(--muted);
  font-size: 18px;
  font-weight: 760;
  text-transform: uppercase;
  letter-spacing: 0.08em;
}
.eyebrow.blue { color: var(--blue); }
.eyebrow.red { color: var(--red); }
.eyebrow.teal { color: var(--teal); }
.lead {
  color: var(--ink-soft);
  font-size: 30px;
  line-height: 1.35;
  margin: 0;
}
.caption {
  color: var(--muted);
  font-size: 22px;
  line-height: 1.35;
}
.big-statement {
  margin-top: 44px;
  font-size: 58px;
  line-height: 1.08;
  font-weight: 800;
  color: var(--red);
}

/* === 7. Layout grids === */
.grid-2 {
  display: grid;
  grid-template-columns: 0.98fr 1.02fr;
  gap: 72px;
  align-items: center;
  height: 760px;
}
.grid-2.narrow-left { grid-template-columns: 0.86fr 1.14fr; }
.columns {
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 38px;
  margin-top: 44px;
}

/* === 8. Panel surface === */
.panel {
  background:
    linear-gradient(180deg, rgba(255,255,255,0.92), rgba(255,255,255,0.76)),
    var(--panel);
  border: 1.5px solid rgba(148, 160, 178, 0.42);
  border-radius: 8px;
  box-shadow: var(--shadow);
  backdrop-filter: blur(2px);
}

/* === 9. Bullet list === */
.bullets {
  display: grid;
  gap: 34px;
  margin-top: 16px;
}
.bullet {
  display: grid;
  grid-template-columns: 24px 1fr;
  gap: 22px;
  align-items: start;
}
.bullet i {
  width: 18px;
  height: 18px;
  margin-top: 12px;
  background: var(--blue);
  border-radius: 5px;
  display: block;
  box-shadow: 0 0 0 7px rgba(21,95,159,0.10);
}
.bullet.red i  { background: var(--red); }
.bullet.teal i { background: var(--teal); }
.bullet strong {
  display: block;
  font-size: 35px;
  line-height: 1.12;
  margin-bottom: 8px;
}
.bullet span {
  display: block;
  color: var(--ink-soft);
  font-size: 24px;
  line-height: 1.24;
}

/* === 10. Tile grid & tiles === */
.tile-grid {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 28px;
  margin-top: 30px;
}
.tile-grid.two {
  grid-template-columns: repeat(2, 1fr);
  width: 1110px;
  margin-left: auto;
  margin-right: auto;
}
.tile {
  position: relative;
  min-height: 240px;
  padding: 34px 34px 30px;
  display: grid;
  grid-template-columns: 108px 1fr;
  gap: 26px;
  align-items: center;
  overflow: hidden;
}
.tile::before {
  content: "";
  position: absolute;
  left: 0;
  top: 0;
  bottom: 0;
  width: 6px;
  background: linear-gradient(180deg, var(--blue), rgba(21,95,159,0.18));
}
.tile:nth-child(2)::before { background: linear-gradient(180deg, var(--red), rgba(199,71,53,0.18)); }
.tile:nth-child(3)::before { background: linear-gradient(180deg, var(--teal), rgba(20,141,128,0.18)); }
.tile-grid.two .tile:first-child::before { background: linear-gradient(180deg, var(--teal), rgba(20,141,128,0.18)); }
.tile-grid.two .tile:nth-child(2)::before { background: linear-gradient(180deg, var(--red), rgba(199,71,53,0.18)); }
.tile > * {
  position: relative;
  z-index: 1;
}
.icon {
  position: relative;
  width: 108px;
  height: 108px;
  border-radius: 8px;
  background: linear-gradient(180deg, #eef6fd, var(--blue-soft));
  color: var(--blue);
  display: grid;
  place-items: center;
  font-size: 48px;
  font-weight: 760;
  box-shadow: inset 0 0 0 1px rgba(21,95,159,0.08), var(--shadow-soft);
}
.icon svg {
  width: 72px;
  height: 72px;
  stroke: currentColor;
  fill: none;
  stroke-width: 4;
  stroke-linecap: round;
  stroke-linejoin: round;
}
.icon svg .fill   { fill: currentColor; stroke: none; }
.icon.red  { background: linear-gradient(180deg, #fff0ed, var(--red-soft));  color: var(--red); }
.icon.teal { background: linear-gradient(180deg, #ebfaf6, var(--teal-soft)); color: var(--teal); }
.tile h2 {
  margin: 0 0 10px;
  font-size: 32px;
  line-height: 1.1;
}
.tile p {
  margin: 0;
  color: var(--ink-soft);
  font-size: 21px;
  line-height: 1.25;
}

/* === 11. Tag badge === */
.tag {
  display: inline-flex;
  align-items: center;
  min-height: 42px;
  padding: 0 20px;
  border-radius: 999px;
  background: rgba(255,255,255,0.78);
  border: 1.5px solid var(--rule);
  color: var(--ink-soft);
  font-size: 19px;
  font-weight: 700;
  white-space: nowrap;
}

/* === 12. Diagram primitives === */
.diagram {
  position: relative;
  min-height: 560px;
}
.node {
  position: absolute;
  width: 84px;
  height: 84px;
  border-radius: 50%;
  display: grid;
  place-items: center;
  color: white;
  font-size: 30px;
  font-weight: 760;
  background: var(--blue);
  box-shadow: 0 14px 30px rgba(31, 95, 159, 0.24), inset 0 0 0 1px rgba(255,255,255,0.26);
  z-index: 2;
}
.node.red  { background: var(--red);  box-shadow: 0 10px 26px rgba(199, 68, 50, 0.20); }
.node.teal { background: var(--teal); box-shadow: 0 10px 26px rgba(23, 140, 125, 0.18); }
.node.gray { background: #6b7685; box-shadow: 0 10px 26px rgba(86, 97, 110, 0.22); }
.line {
  position: absolute;
  height: 4px;
  background: var(--rule);
  transform-origin: left center;
  border-radius: 999px;
  z-index: 1;
}
.line.blue { background: rgba(31,95,159,0.34); }
.line.red  { background: rgba(199,68,50,0.42); }
.arrow {
  position: absolute;
  height: 4px;
  background: rgba(82,97,115,0.78);
  transform-origin: left center;
  border-radius: 999px;
}
.arrow::after {
  content: "";
  position: absolute;
  right: -3px;
  top: -8px;
  border-left: 18px solid rgba(82,97,115,0.78);
  border-top: 10px solid transparent;
  border-bottom: 10px solid transparent;
}

/* === 13. Timeline === */
.timeline {
  position: relative;
  height: 590px;
  margin-top: 28px;
}
.timeline svg {
  position: absolute;
  left: 60px;
  top: 230px;
  width: 1600px;
  height: 160px;
  overflow: visible;
}
.timeline-path {
  fill: none;
  stroke: url(#timelineGradient);
  stroke-width: 7;
  stroke-linecap: round;
  stroke-linejoin: round;
  filter: drop-shadow(0 10px 18px rgba(31,95,159,0.12));
}
.milestone {
  position: absolute;
  width: 220px;
  text-align: center;
}
.milestone .dot {
  position: absolute;
  left: 50%;
  width: 32px;
  height: 32px;
  margin-left: -16px;
  border-radius: 50%;
  background: var(--blue);
  box-shadow: 0 0 0 10px rgba(31,95,159,0.10), 0 10px 24px rgba(31,95,159,0.18);
}
.milestone .dot { top: var(--dot-top, 250px); }
.milestone.teal .dot {
  background: var(--teal);
  box-shadow: 0 0 0 10px rgba(20,141,128,0.12), 0 10px 24px rgba(20,141,128,0.20);
}
.milestone h2 {
  margin: 0 0 8px;
  color: var(--blue);
  font-size: 22px;
}
.milestone.teal h2 { color: var(--teal); }
.milestone strong {
  display: block;
  font-size: 28px;
  margin-bottom: 10px;
}
.milestone p {
  margin: 0;
  color: var(--ink-soft);
  font-size: 18px;
  line-height: 1.25;
}

/* === 14. Policy box & opponents === */
.policy-box {
  position: absolute;
  background: linear-gradient(180deg, #fff, #f7faff);
  border: 4px solid rgba(23,33,44,0.88);
  border-radius: 8px;
  padding: 28px;
  font-size: 33px;
  font-weight: 760;
  line-height: 1.14;
  text-align: center;
  box-shadow: var(--shadow);
}
.opponent {
  position: absolute;
  width: 450px;
  height: 116px;
  border-radius: 8px;
  border: 4px solid var(--blue);
  display: grid;
  place-items: center;
  color: var(--blue);
  background: rgba(255,255,255,0.88);
  font-size: 28px;
  font-weight: 760;
  box-shadow: var(--shadow);
}
.opponent.red   { border-color: var(--red);   color: var(--red); }
.opponent.amber { border-color: var(--amber); color: var(--amber); }

.limitation-diagram {
  height: 700px;
  display: grid;
  grid-template-rows: auto 1fr auto;
  gap: 18px;
}
.limitation-stage {
  position: relative;
  min-height: 430px;
  display: grid;
  place-items: center;
}
.limitation-svg {
  width: 100%;
  height: 100%;
  max-height: 540px;
  display: block;
}

/* === 15. Numbered list === */
.numbered {
  display: grid;
  gap: 44px;
  margin-top: 44px;
}
.numbered .item {
  display: grid;
  grid-template-columns: 92px 1fr;
  gap: 28px;
  align-items: center;
}
.num {
  width: 86px;
  height: 86px;
  border-radius: 50%;
  display: grid;
  place-items: center;
  color: white;
  font-size: 38px;
  font-weight: 800;
  background: var(--blue);
  box-shadow: 0 12px 28px rgba(31, 95, 159, 0.18);
}
.num.red  { background: var(--red); }
.num.teal { background: var(--teal); }
.numbered strong {
  display: block;
  font-size: 42px;
  line-height: 1.08;
}
.numbered span { color: var(--ink-soft); font-size: 25px; }

/* === 16. Process cards === */
.process {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 28px;
  align-items: stretch;
  margin-top: 52px;
}
.process-card {
  position: relative;
  padding: 34px 32px;
  min-height: 210px;
  text-align: center;
  overflow: hidden;
}
.process-card::before {
  content: "";
  position: absolute;
  left: 28px;
  right: 28px;
  top: 0;
  height: 5px;
  border-radius: 0 0 999px 999px;
  background: linear-gradient(90deg, var(--blue), var(--teal));
}
.process-card h2 {
  margin: 0 0 12px;
  font-size: 31px;
}
.process-card p {
  margin: 0;
  color: var(--ink-soft);
  font-size: 21px;
  line-height: 1.28;
}

/* === 17. Slider === */
.slider {
  height: 360px;
  padding: 64px 54px;
}
.track {
  position: relative;
  height: 28px;
  margin: 90px 0 34px;
  border-radius: 999px;
  background: linear-gradient(90deg, var(--red), var(--red-soft) 34%, var(--blue-soft) 66%, var(--blue));
  box-shadow: inset 0 0 0 1px rgba(255,255,255,0.45), var(--shadow-soft);
}
.track::before {
  content: "";
  position: absolute;
  left: 35%;
  top: -26px;
  width: 42px;
  height: 80px;
  border-radius: 8px;
  background: linear-gradient(180deg, #2f3947, var(--ink));
  box-shadow: 0 14px 28px rgba(32,40,50,0.22);
}
.track-labels {
  display: flex;
  justify-content: space-between;
  font-size: 29px;
  font-weight: 760;
}

/* === 18. Note === */
.note {
  padding: 34px 38px;
  font-size: 25px;
  line-height: 1.32;
  color: var(--ink-soft);
  background:
    linear-gradient(90deg, rgba(21,95,159,0.06), transparent 34%),
    rgba(255,255,255,0.82);
}
.note-blue { border-color: var(--blue); }
.note-red { border-color: var(--red); }

/* === 19. Stack / layer === */
.stack {
  width: 620px;
  margin: 46px auto 0;
  display: grid;
  gap: 16px;
}
.layer {
  min-height: 112px;
  padding: 22px 30px;
  border: 3px solid var(--teal);
  background: rgba(255,255,255,0.9);
}
.layer:nth-child(odd) {
  background: linear-gradient(180deg, #effaf7, var(--teal-soft));
}
.layer strong {
  display: block;
  font-size: 29px;
  margin-bottom: 8px;
}
.layer span { color: var(--teal); font-size: 22px; font-weight: 760; }

/* === 20. Outcome panels === */
.outcome {
  position: relative;
  padding: 42px 44px;
  min-height: 548px;
  overflow: hidden;
}
.outcome::before {
  content: "";
  position: absolute;
  left: 0;
  right: 0;
  top: 0;
  height: 7px;
  background: linear-gradient(90deg, var(--blue), var(--teal));
}
.outcome.light::before {
  background: linear-gradient(90deg, var(--teal), rgba(20,141,128,0.25));
}
.outcome > * {
  position: relative;
  z-index: 1;
}
.outcome h2 {
  margin: 0 0 8px;
  color: var(--blue);
  font-size: 22px;
  letter-spacing: 0.08em;
}
.outcome h3 {
  margin: 0 0 40px;
  font-size: 38px;
  line-height: 1.15;
}
.outcome ul {
  margin: 0;
  padding-left: 32px;
  color: var(--ink);
  font-size: 28px;
  line-height: 1.58;
}
.outcome.light              { border-color: var(--teal); }
.outcome.light h2           { color: var(--teal); }
.outcome.light h3,
.outcome.light ul           { color: var(--ink); }

/* === 21. Illustration placeholders === */
.illustration-placeholder {
  position: relative;
  min-height: 520px;
  display: grid;
  place-items: center;
  overflow: hidden;
  padding: 52px;
  text-align: center;
  background:
    linear-gradient(135deg, rgba(21,95,159,0.10), transparent 34%),
    linear-gradient(315deg, rgba(20,141,128,0.10), transparent 34%),
    rgba(255,255,255,0.78);
}
.illustration-placeholder::before {
  content: "";
  position: absolute;
  inset: 28px;
  border: 2px dashed rgba(21,95,159,0.30);
  border-radius: 8px;
}
.illustration-placeholder .placeholder-label {
  display: block;
  color: var(--blue);
  font-size: 52px;
  font-weight: 800;
  letter-spacing: 0.02em;
}
.illustration-placeholder .placeholder-meta {
  display: block;
  max-width: 680px;
  margin-top: 18px;
  color: var(--ink-soft);
  font-size: 24px;
  line-height: 1.28;
}
.title-illustration {
  width: 100%;
  min-height: 0;
  aspect-ratio: 1 / 1;
  place-self: center;
}
.wide-illustration {
  min-height: 610px;
  margin-top: 34px;
}
.side-illustration {
  width: 100%;
  min-height: 0;
  aspect-ratio: 1566 / 1005;
}
.testbed-illustration {
  min-height: 460px;
}
.illustration-image-panel {
  padding: 0;
  background:
    linear-gradient(135deg, rgba(21,95,159,0.06), transparent 42%),
    linear-gradient(315deg, rgba(20,141,128,0.06), transparent 36%),
    rgba(255,255,255,0.72);
}
.illustration-image-panel::before {
  display: none;
}
.title-illustration-image {
  width: 100%;
  height: 100%;
  object-fit: cover;
  display: block;
}
.title-illustration-svg {
  width: 100%;
  height: 100%;
  display: block;
  padding: 36px 28px;
}

/* === 22. Layout-specific polish === */
.split-environment {
  display: grid;
  grid-template-columns: repeat(2, minmax(0, 1fr));
  gap: 28px;
  height: 620px;
}
.split-card {
  position: relative;
  padding: 34px;
  min-height: 100%;
}
.split-card h2 {
  margin: 0 0 22px;
  font-size: 28px;
}
.state-strip {
  display: grid;
  gap: 18px;
  margin-top: 34px;
}
.state-pill {
  min-height: 58px;
  display: flex;
  align-items: center;
  gap: 14px;
  padding: 0 20px;
  border-radius: 8px;
  background: rgba(255,255,255,0.72);
  border: 1.5px solid rgba(148,160,178,0.36);
  color: var(--ink-soft);
  font-size: 21px;
  font-weight: 700;
}
.state-pill i {
  width: 16px;
  height: 16px;
  border-radius: 50%;
  background: var(--blue);
}
.state-pill.red i { background: var(--red); }
.state-pill.teal i { background: var(--teal); }
.scale-tag {
  display: inline-flex;
  min-height: 34px;
  align-items: center;
  padding: 0 13px;
  margin-top: 10px;
  border-radius: 999px;
  background: rgba(21,95,159,0.10);
  color: var(--blue);
  font-size: 15px;
  font-weight: 800;
}
.output-badge {
  position: absolute;
  width: 140px;
  height: 48px;
  display: grid;
  place-items: center;
  border-radius: 999px;
  border: 1.5px solid var(--rule);
  background: rgba(255,255,255,0.86);
  color: var(--muted);
  font-size: 18px;
  font-weight: 760;
}
.metric-grid {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 22px;
  margin-top: 46px;
}
.metric-caption {
  text-align: center;
  margin-top: 30px;
}
.metric-card {
  min-height: 430px;
  padding: 34px 28px;
  text-align: center;
}
.metric-card h2 {
  margin: 0 0 14px;
  font-size: 28px;
  min-height: 32px;
  line-height: 1.1;
}
.metric-card p {
  margin: 0;
  color: var(--ink-soft);
  font-size: 20px;
  line-height: 1.25;
}
.gauge {
  width: 150px;
  height: 150px;
  margin: 34px auto 28px;
  border-radius: 50%;
  display: grid;
  place-items: center;
  color: var(--teal);
  font-size: 34px;
  font-weight: 800;
  background:
    radial-gradient(circle at center, white 0 48%, transparent 49%),
    conic-gradient(var(--teal) 0 68%, rgba(20,141,128,0.16) 68% 100%);
  box-shadow: var(--shadow-soft);
}
.cycle-mark {
  display: grid;
  place-items: center;
  width: 150px;
  height: 150px;
  margin: 34px auto 28px;
  border-radius: 50%;
  background: linear-gradient(135deg, var(--blue-soft), var(--teal-soft));
  color: var(--blue);
  font-size: 30px;
  font-weight: 800;
  box-shadow: var(--shadow-soft);
}
.outcome.primary {
  min-height: 578px;
  border-color: var(--blue) !important;
  border-width: 4px !important;
  box-shadow: 0 30px 88px rgba(21,95,159,0.14);
}

/* === 24. Phase grid (study path) === */
.phase-a { --phase-bg: #dbeafe; --phase-border: #2563eb; --phase-text: #1e40af; }
.phase-b { --phase-bg: #fae8ff; --phase-border: #d946ef; --phase-text: #86198f; }
.phase-c { --phase-bg: #fef3c7; --phase-border: #d97706; --phase-text: #92400e; }
.phase-d { --phase-bg: #ffe4e6; --phase-border: #f43f5e; --phase-text: #be123c; }
.phase-e { --phase-bg: #dcfce7; --phase-border: #059669; --phase-text: #065f46; }
.phase-f { --phase-bg: #cffafe; --phase-border: #0891b2; --phase-text: #155e75; }
.phase-g { --phase-bg: #ede9fe; --phase-border: #7c3aed; --phase-text: #4c1d95; }

.phase-lead {
  margin: 6px 0 30px;
}
.phase-grid {
  display: grid;
  grid-template-columns: repeat(12, 1fr);
  gap: 24px;
}
.phase-card {
  grid-column: span 3;
  padding: 26px 28px 24px;
  display: flex;
  flex-direction: column;
  min-height: 250px;
  background:
    linear-gradient(180deg, var(--phase-bg, #ffffff), rgba(255,255,255,0.65));
  border: 2px solid var(--phase-border, var(--rule));
}
.phase-card.wide {
  grid-column: span 4;
}
.phase-badge {
  width: 54px;
  height: 54px;
  border-radius: 12px;
  background: var(--phase-border, var(--blue));
  color: white;
  display: grid;
  place-items: center;
  font-size: 28px;
  font-weight: 800;
  font-family: Aptos, "Segoe UI", Arial, sans-serif;
  box-shadow: 0 10px 22px rgba(0, 0, 0, 0.12);
}
.phase-card h2 {
  margin: 16px 0 4px;
  font-size: 24px;
  line-height: 1.14;
  color: var(--phase-text, var(--ink));
}
.phase-meta {
  display: block;
  font-size: 15px;
  color: var(--phase-border, var(--muted));
  font-weight: 760;
  text-transform: uppercase;
  letter-spacing: 0.06em;
  opacity: 0.85;
}
.phase-card p {
  margin: 14px 0 0;
  font-size: 17px;
  color: var(--ink-soft);
  line-height: 1.34;
}
.phase-footer {
  text-align: center;
  margin-top: 26px;
}

/* === 23. Responsive-friendly slide compositions === */
.title-composition {
  height: 100%;
  display: grid;
  grid-template-columns: minmax(0, 1.08fr) minmax(560px, 0.82fr);
  gap: 72px;
  align-items: center;
}
.title-copy {
  display: flex;
  flex-direction: column;
  align-items: flex-start;
}
.title-logo {
  width: 144px;
  height: auto;
  margin-bottom: 74px;
}
.title-heading {
  margin: 0 0 34px;
  font-size: clamp(62px, 3.7vw, 70px);
  line-height: 1.02;
  max-width: 980px;
}
.title-subtitle {
  margin: 0;
  color: var(--blue);
  font-size: 34px;
  line-height: 1.18;
  font-style: italic;
}
.title-copy .title-rule {
  margin-top: 42px;
}
.title-author {
  margin: 52px 0 0;
  font-size: 31px;
  font-weight: 760;
}
.title-meta {
  margin: 12px 0 0;
  color: var(--ink-soft);
  font-size: 25px;
}
.split-card {
  display: flex;
  flex-direction: column;
}
.split-card .state-strip {
  margin-top: auto;
}
.agent-svg-wrap {
  flex: 1;
  min-height: 260px;
  display: grid;
  place-items: center;
  padding: 12px 0;
}
.agent-svg {
  width: 100%;
  height: 100%;
  max-height: 320px;
  display: block;
}
.domain-row {
  display: grid;
  grid-template-columns: repeat(5, minmax(0, 1fr));
  gap: 14px;
  margin-top: 22px;
}
.domain-chip {
  min-height: 86px;
  padding: 18px;
  display: grid;
  place-items: center;
  text-align: center;
}
.domain-chip h2 {
  margin: 0;
  font-size: 20px;
}
.testbed-row {
  display: grid;
  grid-template-columns: repeat(4, minmax(0, 1fr));
  gap: 18px;
  margin-top: 20px;
}
.testbed-chip {
  min-height: 86px;
  padding: 16px;
  text-align: center;
  display: grid;
  align-content: center;
}
.testbed-chip h2 {
  margin: 0;
  font-size: 20px;
}
.testbed-chip p {
  margin: 5px 0 0;
  font-size: 16px;
}
.slide-bottom-note {
  position: absolute;
  left: 104px;
  right: 104px;
  bottom: 104px;
  text-align: center;
}
.contribution-layout {
  margin-top: -30px;
}
.visual-with-note {
  display: flex;
  flex-direction: column;
  gap: 24px;
}
.visual-with-note .note {
  margin-top: 0;
}
.closing-layout {
  display: flex;
  flex-direction: column;
  gap: 22px;
}
.closing-illustration {
  width: 100%;
  min-height: 0;
  aspect-ratio: 2172 / 724;
  margin-top: 0;
}
.closing-summary {
  display: grid;
  grid-template-columns: repeat(2, minmax(0, 1fr));
  gap: 32px;
}
.summary-panel {
  padding: 22px 28px;
}
.summary-panel strong {
  display: block;
  font-size: 24px;
  color: var(--muted);
}
.summary-panel span {
  display: block;
  margin-top: 8px;
  font-size: 25px;
  color: var(--ink-soft);
}
.summary-panel.next {
  border-color: var(--blue);
}
.summary-panel.next strong {
  color: var(--blue);
}
.summary-panel.next span {
  color: var(--ink);
}
.thank-you {
  margin: 18px 0 0;
  padding: 24px 0;
  text-align: center;
  font-size: 40px;
  font-weight: 800;
  color: var(--blue);
}
.timeline-lead {
  margin-top: 8px;
}
.limitation-lead {
  margin-top: 20px;
}
.diagram-caption-bottom {
  position: absolute;
  left: 0;
  right: 0;
  bottom: 36px;
  text-align: center;
}
.diagram-caption-low {
  position: absolute;
  left: 0;
  right: 0;
  bottom: 110px;
}
.problem-title-blue { color: var(--blue); }
.problem-title-red { color: var(--red); }
.problem-title-teal { color: var(--teal); }
.open-problem-visual {
  width: min(650px, 100%);
  aspect-ratio: 1 / 1;
  place-self: center;
}
.open-problem-image {
  width: 100%;
  height: 100%;
  object-fit: cover;
  display: block;
}

.example-grid {
  display: grid;
  gap: 22px;
  align-items: stretch;
}
.example-grid.five {
  grid-template-columns: repeat(5, minmax(0, 1fr));
  margin-top: 42px;
}
.example-grid.four {
  grid-template-columns: repeat(4, minmax(0, 1fr));
  margin-top: 36px;
}
.example-card {
  min-height: 570px;
  padding: 0;
  display: flex;
  flex-direction: column;
  overflow: hidden;
}
.example-grid.four .example-card {
  min-height: 555px;
}
.example-thumb-frame {
  height: 250px;
  margin: 18px 18px 0;
  border-radius: 8px;
  overflow: hidden;
  background: linear-gradient(135deg, var(--blue-soft), var(--teal-soft));
  border: 1.5px solid rgba(148,160,178,0.30);
}
.example-grid.four .example-thumb-frame {
  height: 270px;
}
.example-thumb {
  display: block;
  width: 100%;
  height: 100%;
  object-fit: cover;
  filter: saturate(0.86) contrast(1.05);
}
.placeholder-thumb {
  display: grid;
  place-items: center;
  padding: 18px;
  color: var(--blue);
  font-size: 24px;
  font-weight: 800;
  text-align: center;
}
.example-card-content {
  padding: 24px 24px 28px;
}
.example-card h2 {
  margin: 0 0 14px;
  font-size: 28px;
  line-height: 1.08;
}
.example-bullets {
  margin: 0;
  padding: 0;
  list-style: none;
  display: grid;
  gap: 10px;
}
.example-bullets li {
  position: relative;
  padding-left: 22px;
  color: var(--ink-soft);
  font-size: 20px;
  line-height: 1.28;
}
.example-bullets li::before {
  content: "";
  position: absolute;
  left: 0;
  top: 9px;
  width: 10px;
  height: 10px;
  border-radius: 3px;
  background: var(--blue);
}
"""


def write_common_css() -> Path:
    HTML_DIR.mkdir(parents=True, exist_ok=True)
    path = HTML_DIR / "common.css"
    path.write_text(COMMON_CSS, encoding="utf-8")
    return path


def shell(title: str, body: str, *, slide_no: int | None, accent: str = "blue", title_slide: bool = False) -> str:
    cls = "slide title-slide" if title_slide else "slide"
    header = "" if title_slide else f"""
      <header class="slide-header">
        <h1>{esc(title)}</h1>
        <div class="title-rule {accent}"></div>
      </header>
    """
    foot = "" if slide_no is None else footer(slide_no, title)
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width={WIDTH}, initial-scale=1">
  <title>{esc(title)}</title>
  <link rel="stylesheet" href="common.css">
</head>
<body>
  <main class="{cls}">
    {header}
    {body}
    {foot}
  </main>
</body>
</html>
"""


def line(x1: int, y1: int, x2: int, y2: int, cls: str = "") -> str:
    import math

    length = math.hypot(x2 - x1, y2 - y1)
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    return (
        f'<div class="line {cls}" style="left:{x1}px;top:{y1}px;'
        f'width:{length:.1f}px;transform:rotate({angle:.2f}deg)"></div>'
    )


def arrow(x1: int, y1: int, x2: int, y2: int) -> str:
    import math

    length = math.hypot(x2 - x1, y2 - y1)
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    return (
        f'<div class="arrow" style="left:{x1}px;top:{y1}px;'
        f'width:{length:.1f}px;transform:rotate({angle:.2f}deg)"></div>'
    )


def svg_icon(kind: str) -> str:
    icons = {
        "market": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <path d="M14 64V18M14 64H68"/>
            <path d="M20 54L32 43L42 49L54 30L66 35"/>
            <path d="M58 30H68V40"/>
          </svg>
        """,
        "lock": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <rect x="18" y="34" width="44" height="32" rx="6"/>
            <path d="M28 34V25C28 15 35 9 40 9C45 9 52 15 52 25V34"/>
            <circle class="fill" cx="40" cy="50" r="4"/>
          </svg>
        """,
        "network": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <path d="M40 40L19 22M40 40L62 23M40 40L22 61M40 40L60 61"/>
            <circle cx="40" cy="40" r="8"/>
            <circle cx="19" cy="22" r="7"/>
            <circle cx="62" cy="23" r="7"/>
            <circle cx="22" cy="61" r="7"/>
            <circle cx="60" cy="61" r="7"/>
          </svg>
        """,
        "grid": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <path d="M16 16H64V64H16zM32 16V64M48 16V64M16 32H64M16 48H64"/>
            <circle class="fill" cx="24" cy="24" r="5"/>
            <circle class="fill" cx="56" cy="40" r="5"/>
          </svg>
        """,
        "cards": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <rect x="18" y="18" width="30" height="42" rx="5"/>
            <rect x="32" y="24" width="30" height="42" rx="5"/>
            <path d="M47 38C42 32 36 37 40 43C43 47 47 51 47 51C47 51 51 47 54 43C58 37 52 32 47 38Z"/>
          </svg>
        """,
        "gavel": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <path d="M18 58H54"/>
            <path d="M28 24L44 40M22 30L34 18M38 46L50 34"/>
            <path d="M46 42L64 60"/>
          </svg>
        """,
        "coalition": """
          <svg viewBox="0 0 80 80" aria-hidden="true">
            <path d="M22 24H58M22 24L18 58M58 24L62 58M18 58H62M22 24L62 58"/>
            <circle class="fill" cx="22" cy="24" r="6"/>
            <circle class="fill" cx="58" cy="24" r="6"/>
            <circle class="fill" cx="18" cy="58" r="6"/>
            <circle class="fill" cx="62" cy="58" r="6"/>
          </svg>
        """,
    }
    return icons[kind]


def slides() -> list[tuple[str, str]]:
    title_visual = asset_img(
        "illustration1.png",
        "Multi-agent imperfect-information game-tree hero illustration",
        "title-illustration-image",
    )
    if title_visual:
        title_visual_html = f'<div class="title-illustration panel illustration-image-panel">{title_visual}</div>'
    else:
        title_visual_html = """
          <div class="illustration-placeholder title-illustration panel">
            <div>
              <span class="placeholder-label">illustration1</span>
              <span class="placeholder-meta">Hero visual: multi-agent network, hidden information, game-tree structure</span>
            </div>
          </div>
        """

    s1 = shell(
        "Adaptive Strategy Learning",
        f"""
        <section class="title-composition">
          <div class="title-copy">
            {'<img class="title-logo" src="' + LOGO_URI + '" alt="Ruse University logo">' if LOGO_URI else ''}
            <h1 class="title-heading">
              Adaptive Strategy Learning<br>
              in Multi-Agent<br>
              Imperfect-Information Games
            </h1>
            <p class="title-subtitle">
              From Equilibrium Computation to Safe Opponent Exploitation
            </p>
            <div class="title-rule blue"></div>
            <p class="title-author">Alexander Andreev<br>Prof. Dr. Tsvetomir Vasilev</p>
            <p class="title-meta">
              PhD session report &middot; May 2026 &middot; Ruse University "Angel Kanchev"
            </p>
          </div>
          {title_visual_html}
        </section>
        """,
        slide_no=None,
        title_slide=True,
    )

    single_agent_svg = """
        <svg class="agent-svg" viewBox="0 0 400 360" preserveAspectRatio="xMidYMid meet" aria-hidden="true">
          <circle cx="200" cy="180" r="58" fill="#155f9f"/>
          <text x="200" y="201" text-anchor="middle" fill="white" font-size="56" font-weight="760"
                font-family="Aptos, Segoe UI, Arial, sans-serif">A</text>
        </svg>
    """
    multi_agent_svg = """
        <svg class="agent-svg" viewBox="0 0 400 360" preserveAspectRatio="xMidYMid meet" aria-hidden="true">
          <g stroke="#c74735" stroke-width="3" stroke-linecap="round" opacity="0.55" fill="none">
            <line x1="200" y1="180" x2="100" y2="80"/>
            <line x1="200" y1="180" x2="300" y2="80"/>
            <line x1="200" y1="180" x2="60"  y2="240"/>
            <line x1="200" y1="180" x2="340" y2="240"/>
            <line x1="200" y1="180" x2="200" y2="320"/>
            <line x1="100" y1="80"  x2="300" y2="80"/>
            <line x1="60"  y1="240" x2="200" y2="320"/>
            <line x1="340" y1="240" x2="200" y2="320"/>
          </g>
          <circle cx="200" cy="180" r="42" fill="#155f9f"/>
          <text x="200" y="195" text-anchor="middle" fill="white" font-size="34" font-weight="760"
                font-family="Aptos, Segoe UI, Arial, sans-serif">A</text>
          <circle cx="100" cy="80"  r="30" fill="#c74735"/>
          <text   x="100" y="92"  text-anchor="middle" fill="white" font-size="26" font-weight="760">?</text>
          <circle cx="300" cy="80"  r="30" fill="#c74735"/>
          <text   x="300" y="92"  text-anchor="middle" fill="white" font-size="26" font-weight="760">?</text>
          <circle cx="60"  cy="240" r="30" fill="#c74735"/>
          <text   x="60"  y="252" text-anchor="middle" fill="white" font-size="26" font-weight="760">?</text>
          <circle cx="340" cy="240" r="30" fill="#c74735"/>
          <text   x="340" y="252" text-anchor="middle" fill="white" font-size="26" font-weight="760">?</text>
          <circle cx="200" cy="320" r="30" fill="#c74735"/>
          <text   x="200" y="332" text-anchor="middle" fill="white" font-size="26" font-weight="760">?</text>
        </svg>
    """
    s2 = shell(
        "Why this matters now",
        f"""
        <section class="grid-2">
          <div class="bullets">
            <div class="bullet"><i></i><div><strong>Strong single agents</strong><span>Well understood in constrained benchmarks</span></div></div>
            <div class="bullet red"><i></i><div><strong>Agents interacting with other agents</strong><span>Open problem</span></div></div>
            <div class="bullet red"><i></i><div><strong>Hidden information</strong><span>Adversarial intent &middot; Real-time decisions</span></div></div>
          </div>
          <div class="split-environment">
            <article class="split-card panel">
              <h2>Benchmark agent</h2>
              <div class="agent-svg-wrap">{single_agent_svg}</div>
              <div class="state-strip">
                <div class="state-pill"><i></i>Single objective</div>
                <div class="state-pill"><i></i>Static environment</div>
                <div class="state-pill"><i></i>Solo evaluation</div>
              </div>
            </article>
            <article class="split-card panel">
              <h2>Strategic environment</h2>
              <div class="agent-svg-wrap">{multi_agent_svg}</div>
              <div class="state-strip">
                <div class="state-pill red"><i></i>Hidden intent</div>
                <div class="state-pill red"><i></i>Adversarial response</div>
                <div class="state-pill red"><i></i>Real-time pressure</div>
              </div>
            </article>
          </div>
        </section>
        """,
        slide_no=2,
        accent="blue",
    )

    use_tiles = [
        ("Financial markets", "Hidden positions &middot; Microsecond reactions", svg_icon("market"), ""),
        ("Cybersecurity", "Adaptive attackers &middot; Hidden intent", svg_icon("lock"), "red"),
        ("Social platforms", "Bot networks &middot; Coordinated disinformation", svg_icon("network"), ""),
        ("Security", "Patrol and screening randomisation", svg_icon("grid"), "teal"),
        ("Gaming platforms", "Collusion &middot; Fraud detection", svg_icon("cards"), "red"),
    ]
    tile_html = "".join(
        f"""
        <article class="tile panel">
          <div class="icon {color}">{icon}</div>
          <div><h2>{title}</h2><p>{sub}</p></div>
        </article>
        """
        for title, sub, icon, color in use_tiles[:3]
    )
    tile_html_2 = "".join(
        f"""
        <article class="tile panel">
          <div class="icon {color}">{icon}</div>
          <div><h2>{title}</h2><p>{sub}</p></div>
        </article>
        """
        for title, sub, icon, color in use_tiles[3:]
    )
    use_cases_data = [
        ("Financial markets", "example_thumbs/financial_markets.jpg",
         ["Hidden positions", "Microsecond reactions", "Adversarial flow"]),
        ("Cybersecurity", "example_thumbs/cybersecurity.jpg",
         ["Adaptive attackers", "Hidden intent", "Real-time defense"]),
        ("Social platforms", "example_thumbs/social_platforms.jpeg",
         ["Bot networks", "Coordinated disinformation", "Mixed human / agent traffic"]),
        ("Security", "example_thumbs/security.webp",
         ["Patrol randomisation", "Screening allocation", "Deployed today"]),
        ("Gaming platforms", "example_thumbs/gaming_platforms.jpg",
         ["Collusion detection", "Fraud signals", "Online play at scale"]),
    ]
    use_card_html = "".join(
        f"""
        <article class="example-card panel">
          <div class="example-thumb-frame">{thumbnail(image, title)}</div>
          <div class="example-card-content">
            <h2>{title}</h2>
            <ul class="example-bullets">
              {''.join(f'<li>{b}</li>' for b in bullets)}
            </ul>
          </div>
        </article>
        """
        for title, image, bullets in use_cases_data
    )
    s3 = shell(
        "Where this shows up",
        f"""
        <section class="example-grid five">
          {use_card_html}
        </section>
        """,
        slide_no=3,
        accent="blue",
    )

    testbeds = [
        ("Poker / Belot-like", "Hidden cards &middot; Inference", svg_icon("cards")),
        ("Auction-style", "Hidden valuations &middot; Bidding", svg_icon("gavel")),
        ("Pursuit-evasion", "Partial observability &middot; Adversarial search", svg_icon("grid")),
        ("Coalition games", "Alliances &middot; Betrayal &middot; N-player incentives", svg_icon("coalition")),
    ]
    testbed_data = [
        ("Poker / Belot-like", "example_thumbs/poker_belot.jpg",
         ["Hidden cards", "Opponent inference", "Sequential betting"]),
        ("Auction-style", "example_thumbs/auction.webp",
         ["Private valuations", "Strategic bidding", "One-shot decisions"]),
        ("Pursuit-evasion", "example_thumbs/pursuit_evasion.png",
         ["Partial observability", "Adversarial search", "Spatial inference"]),
        ("Coalition games", "example_thumbs/coalition_diplomacy.jpeg",
         ["Alliances and betrayal", "N-player incentives", "Negotiation pressure"]),
    ]
    testbed_card_html = "".join(
        f"""
        <article class="example-card panel">
          <div class="example-thumb-frame">{thumbnail(image, title)}</div>
          <div class="example-card-content">
            <h2>{title}</h2>
            <ul class="example-bullets">
              {''.join(f'<li>{b}</li>' for b in bullets)}
            </ul>
          </div>
        </article>
        """
        for title, image, bullets in testbed_data
    )
    s4 = shell(
        "Why study games?",
        f"""
        <p class="lead">Controlled environments where strategic behaviour is testable.</p>
        <section class="example-grid four">
          {testbed_card_html}
        </section>
        <p class="caption slide-bottom-note">
          Validation environments, not the final application domains.
        </p>
        """,
        slide_no=4,
        accent="blue",
    )

    milestones = [
        ("2007", "CFR",                "Regret minimisation<br>imperfect-info games",       "Algorithmic base",  70,  "top",    280, "blue"),
        ("2016", "AlphaGo",             "Deep RL + MCTS<br>Go",                              "Perfect info",       295, "bottom", -46, "teal"),
        ("2017", "Libratus",            "2-Player poker<br>Superhuman",                      "2 players",          520, "top",    280, "blue"),
        ("2017", "AlphaZero",           "Self-play<br>Chess &middot; Shogi &middot; Go",     "Perfect info",       745, "bottom", -46, "teal"),
        ("2019", "Pluribus",            "6-Player poker<br>Superhuman",                      "6 players",          970, "top",    280, "blue"),
        ("2020", "ReBeL",               "Search + RL<br>imperfect-info",                     "Search + learning",  1195, "bottom", -46, "blue"),
        ("2022", "DeepNash / CICERO",   "Stratego &middot; Diplomacy<br>7 players",          "Large scale",        1420, "top",    280, "blue"),
    ]
    m_html = "".join(
        f"""
        <div class="milestone {pos} {color}" style="left:{x}px;top:{70 if pos == 'top' else 338}px;--dot-top:{dot_top}px">
          <div class="dot"></div><h2>{year}</h2><strong>{name}</strong><p>{sub}</p><span class="scale-tag">{scale}</span>
        </div>
        """
        for year, name, sub, scale, x, pos, dot_top, color in milestones
    )
    s5 = shell(
        "State of the art",
        f"""
        <p class="lead timeline-lead">"Superhuman" in progressively larger games.</p>
        <section class="timeline">
          <svg viewBox="0 0 1600 160" aria-hidden="true">
            <defs>
              <linearGradient id="timelineGradient" x1="0" x2="1" y1="0" y2="0">
                <stop offset="0%"     stop-color="#155f9f" />
                <stop offset="16.67%" stop-color="#148d80" />
                <stop offset="33.33%" stop-color="#155f9f" />
                <stop offset="50%"    stop-color="#148d80" />
                <stop offset="66.67%" stop-color="#155f9f" />
                <stop offset="100%"   stop-color="#155f9f" />
              </linearGradient>
            </defs>
            <path class="timeline-path"
                  d="M 120 136
                     C 232 136, 232 78, 345 78
                     C 457 78, 457 136, 570 136
                     C 682 136, 682 78, 795 78
                     C 907 78, 907 136, 1020 136
                     C 1132 136, 1132 78, 1245 78
                     C 1357 78, 1357 136, 1470 136" />
          </svg>
          {m_html}
        </section>
        """,
        slide_no=5,
        accent="blue",
    )

    limitation_svg = """
        <svg class="limitation-svg" viewBox="0 0 1400 540" preserveAspectRatio="xMidYMid meet" aria-hidden="true">
          <defs>
            <marker id="arrow-head-grey" markerWidth="12" markerHeight="12" refX="10" refY="6" orient="auto">
              <polygon points="0 0, 11 6, 0 12" fill="#526173"/>
            </marker>
            <linearGradient id="lim-policy-bg" x1="0" y1="0" x2="0" y2="1">
              <stop offset="0%" stop-color="#ffffff"/>
              <stop offset="100%" stop-color="#f1f4f9"/>
            </linearGradient>
            <linearGradient id="lim-blue-bg" x1="0" y1="0" x2="1" y2="0">
              <stop offset="0%" stop-color="#f4faff"/>
              <stop offset="100%" stop-color="#ffffff"/>
            </linearGradient>
            <linearGradient id="lim-red-bg" x1="0" y1="0" x2="1" y2="0">
              <stop offset="0%" stop-color="#fff4f1"/>
              <stop offset="100%" stop-color="#ffffff"/>
            </linearGradient>
            <linearGradient id="lim-amber-bg" x1="0" y1="0" x2="1" y2="0">
              <stop offset="0%" stop-color="#fff8eb"/>
              <stop offset="100%" stop-color="#ffffff"/>
            </linearGradient>
          </defs>

          <!-- Equilibrium policy box (left) -->
          <g font-family="Aptos, Segoe UI, Arial, sans-serif">
            <rect x="30" y="180" width="380" height="180" rx="12"
                  fill="url(#lim-policy-bg)" stroke="#17212c" stroke-width="4"/>
            <text x="220" y="218" text-anchor="middle" font-size="18" font-weight="800"
                  fill="#8793a5" letter-spacing="3">EQUILIBRIUM POLICY</text>
            <text x="220" y="298" text-anchor="middle" font-size="92" font-weight="700"
                  font-family="Georgia, 'Times New Roman', serif" fill="#17212c" font-style="italic">&#960;*</text>
            <text x="220" y="338" text-anchor="middle" font-size="21" font-weight="700"
                  fill="#526173">Nash equilibrium &#183; pre-computed</text>
          </g>

          <!-- "Same action, every time" eyebrow -->
          <text x="655" y="60" text-anchor="middle" font-size="18" font-weight="800"
                fill="#8793a5" letter-spacing="3"
                font-family="Aptos, Segoe UI, Arial, sans-serif">SAME ACTION, EVERY TIME</text>

          <!-- Arrows from policy to opponents -->
          <g stroke="#526173" stroke-width="3.2" fill="none" marker-end="url(#arrow-head-grey)" opacity="0.78">
            <line x1="410" y1="248" x2="892" y2="80"/>
            <line x1="410" y1="270" x2="892" y2="270"/>
            <line x1="410" y1="292" x2="892" y2="460"/>
          </g>

          <!-- Identical action tokens riding each arrow -->
          <g font-family="Aptos, Segoe UI, Arial, sans-serif">
            <g transform="translate(645, 138)">
              <rect x="-58" y="-23" width="116" height="46" rx="10"
                    fill="white" stroke="#17212c" stroke-width="2"/>
              <text x="0" y="7" text-anchor="middle" font-size="20" font-weight="800" fill="#17212c">play a*</text>
            </g>
            <g transform="translate(645, 270)">
              <rect x="-58" y="-23" width="116" height="46" rx="10"
                    fill="white" stroke="#17212c" stroke-width="2"/>
              <text x="0" y="7" text-anchor="middle" font-size="20" font-weight="800" fill="#17212c">play a*</text>
            </g>
            <g transform="translate(645, 402)">
              <rect x="-58" y="-23" width="116" height="46" rx="10"
                    fill="white" stroke="#17212c" stroke-width="2"/>
              <text x="0" y="7" text-anchor="middle" font-size="20" font-weight="800" fill="#17212c">play a*</text>
            </g>
          </g>

          <!-- Opponent 1: Cautious beginner -->
          <g font-family="Aptos, Segoe UI, Arial, sans-serif">
            <rect x="900" y="20" width="475" height="120" rx="10"
                  fill="url(#lim-blue-bg)" stroke="#155f9f" stroke-width="3.5"/>
            <text x="1137" y="60"  text-anchor="middle" font-size="28" font-weight="800" fill="#155f9f">Cautious beginner</text>
            <text x="1137" y="90"  text-anchor="middle" font-size="18" font-weight="700" fill="#526173">Folds 70% &#183; undersized bets</text>
            <text x="1137" y="122" text-anchor="middle">
              <tspan font-size="15" font-weight="800" fill="#c74735" letter-spacing="1.5">&#10005; MISSED</tspan>
              <tspan font-size="16" font-weight="600" fill="#8793a5" font-style="italic">  value not extracted</tspan>
            </text>
          </g>

          <!-- Opponent 2: Aggressive bluffer -->
          <g font-family="Aptos, Segoe UI, Arial, sans-serif">
            <rect x="900" y="210" width="475" height="120" rx="10"
                  fill="url(#lim-red-bg)" stroke="#c74735" stroke-width="3.5"/>
            <text x="1137" y="250" text-anchor="middle" font-size="28" font-weight="800" fill="#c74735">Aggressive bluffer</text>
            <text x="1137" y="280" text-anchor="middle" font-size="18" font-weight="700" fill="#526173">Raises 60% &#183; weak hands</text>
            <text x="1137" y="312" text-anchor="middle">
              <tspan font-size="15" font-weight="800" fill="#c74735" letter-spacing="1.5">&#10005; MISSED</tspan>
              <tspan font-size="16" font-weight="600" fill="#8793a5" font-style="italic">  bluffs not punished</tspan>
            </text>
          </g>

          <!-- Opponent 3: Coordinated group -->
          <g font-family="Aptos, Segoe UI, Arial, sans-serif">
            <rect x="900" y="400" width="475" height="120" rx="10"
                  fill="url(#lim-amber-bg)" stroke="#bb7d1c" stroke-width="3.5"/>
            <text x="1137" y="440" text-anchor="middle" font-size="28" font-weight="800" fill="#bb7d1c">Coordinated group</text>
            <text x="1137" y="470" text-anchor="middle" font-size="18" font-weight="700" fill="#526173">Shared signals &#183; collusion</text>
            <text x="1137" y="502" text-anchor="middle">
              <tspan font-size="15" font-weight="800" fill="#c74735" letter-spacing="1.5">&#10005; MISSED</tspan>
              <tspan font-size="16" font-weight="600" fill="#8793a5" font-style="italic">  collusion not detected</tspan>
            </text>
          </g>
        </svg>
    """
    s6 = shell(
        "The limitation",
        f"""
        <section class="limitation-diagram">
          <div>
            <div class="big-statement">Fixed strategies. No adaptation.</div>
            <p class="lead limitation-lead">The same policy against every opponent.</p>
          </div>
          <div class="limitation-stage">{limitation_svg}</div>
          <p class="caption" style="text-align:center;margin:0">Safe, but blind. Systematically misses exploitable weaknesses.</p>
        </section>
        """,
        slide_no=6,
        accent="red",
    )

    # .diagram has position:relative; right panel ≈ 935px wide × 650px tall.
    # Four nodes (160×160) at cardinal positions around centre (467, 290).
    # Arrows connect edge-to-edge in clockwise order.
    illustration2 = asset_img(
        "illustration2.png",
        "Play, study, steer safely, and evaluate research cycle",
        "open-problem-image",
    )
    s7_visual = illustration2 or '<div class="placeholder-thumb open-problem-image">illustration2</div>'
    s7 = shell(
        "Three open problems",
        f"""
        <section class="grid-2 narrow-left">
          <div class="numbered">
            <div class="item"><div class="num">1</div><div><strong class="problem-title-blue">Study</strong><span>Infer opponents in real time</span></div></div>
            <div class="item"><div class="num red">2</div><div><strong class="problem-title-red">Steer Safely</strong><span>Exploit without exposing yourself</span></div></div>
            <div class="item"><div class="num teal">3</div><div><strong class="problem-title-teal">Evaluate</strong><span>Measure if an agent adapts well</span></div></div>
          </div>
          <div class="diagram panel illustration-image-panel open-problem-visual">{s7_visual}</div>
        </section>
        """,
        slide_no=7,
        accent="",
    )

    illustration4 = asset_img(
        "illustration4.png",
        "Belief update diagram: observed actions, hidden state, behavioural type, confidence gate",
        "title-illustration-image",
    )
    if illustration4:
        s8_visual = f'<section class="side-illustration panel illustration-image-panel">{illustration4}</section>'
    else:
        s8_visual = """
            <section class="illustration-placeholder side-illustration panel">
              <div>
                <span class="placeholder-label">illustration4</span>
                <span class="placeholder-meta">Belief update diagram: observed actions, hidden state, behavioural type, confidence gate, equilibrium fallback</span>
              </div>
            </section>
        """
    s8 = shell(
        "Behavioral Adaptation Framework",
        f"""
        <div class="eyebrow blue">Contribution 1</div>
        <section class="grid-2 narrow-left contribution-layout">
          <div class="bullets">
            <div class="bullet"><i></i><div><strong>Infer opponent type</strong><span>Not just hidden state</span></div></div>
            <div class="bullet"><i></i><div><strong>Adapt only when evidence is strong enough</strong><span>Fallback to equilibrium</span></div></div>
            <div class="bullet"><i></i><div><strong>Detect anomalies</strong><span>Bots &middot; Collusion &middot; Adversarial users</span></div></div>
          </div>
          <div class="visual-with-note">
            {s8_visual}
            <div class="note panel note-blue">
              Weak evidence keeps the agent near equilibrium play. Strong evidence opens the adaptation path.
            </div>
          </div>
        </section>
        """,
        slide_no=8,
        accent="blue",
    )

    illustration5 = asset_img(
        "illustration5.png",
        "Reference policy band with constrained exploitative policy path",
        "title-illustration-image",
    )
    if illustration5:
        s9_visual = f'<section class="side-illustration panel illustration-image-panel">{illustration5}</section>'
    else:
        s9_visual = """
            <section class="illustration-placeholder side-illustration panel">
              <div>
                <span class="placeholder-label">illustration5</span>
                <span class="placeholder-meta">Reference policy band with constrained exploitative policy path and pi_KL distance</span>
              </div>
            </section>
        """
    s9 = shell(
        "Multi-Agent Safe Exploitation",
        f"""
        <div class="eyebrow red">Contribution 2</div>
        <section class="grid-2 narrow-left contribution-layout">
          <div class="bullets">
            <div class="bullet red"><i></i><div><strong>Exploit detected weaknesses</strong><span>Act on the inferred type</span></div></div>
            <div class="bullet red"><i></i><div><strong>Stay close to a reference policy</strong><span>pi_KL Regularisation</span></div></div>
            <div class="bullet red"><i></i><div><strong>Useful safety heuristics</strong><span>Beyond two-player guarantees</span></div></div>
          </div>
          <div class="visual-with-note">
            {s9_visual}
            <div class="note panel note-red">
              Two-player safety guarantees provably fail at N >= 3. The target is practical, characterised safety behaviour in larger games.
            </div>
          </div>
        </section>
        """,
        slide_no=9,
        accent="red",
    )

    s10 = shell(
        "Evaluation Methodology",
        """
        <div class="eyebrow teal">Contribution 3</div>
        <section class="grid-2 narrow-left contribution-layout">
          <div class="bullets">
            <div class="bullet teal"><i></i><div><strong>Safety</strong><span>Worst-case vulnerability</span></div></div>
            <div class="bullet teal"><i></i><div><strong>Population ranking</strong><span>Non-transitive dynamics</span></div></div>
            <div class="bullet teal"><i></i><div><strong>Statistical reliability</strong><span>Variance reduction</span></div></div>
          </div>
          <div>
            <section class="metric-grid">
              <article class="metric-card panel">
                <h2>Safety</h2>
                <div class="gauge">!</div>
                <p>How vulnerable is the agent in the worst case across the player count?</p>
              </article>
              <article class="metric-card panel">
                <h2>Ranking</h2>
                <div class="cycle-mark">&#8635;</div>
                <p>Does the population have cyclic, non-transitive structure rather than a single best?</p>
              </article>
              <article class="metric-card panel">
                <h2>Reliability</h2>
                <div class="gauge">&#963;</div>
                <p>Are comparisons stable when hidden information drives high variance?</p>
              </article>
            </section>
            <p class="caption metric-caption">Validated across structurally different game types</p>
          </div>
        </section>
        """,
        slide_no=10,
        accent="teal",
    )

    s11 = shell(
        "Study Plan",
        """
        <p class="lead phase-lead">Seven phases of preparation. Each builds the foundation for the contributions.</p>
        <section class="phase-grid">
          <article class="phase-card panel phase-a">
            <div class="phase-badge">A</div>
            <h2>Foundation</h2>
            <span class="phase-meta">Steps 1–2 &middot; Done</span>
            <p>Reinforcement learning &middot; Counterfactual regret minimisation</p>
          </article>
          <article class="phase-card panel phase-b">
            <div class="phase-badge">B</div>
            <h2>Scaling the toolbox</h2>
            <span class="phase-meta">Steps 3–4 &middot; Done</span>
            <p>CFR variants &middot; Game abstraction &middot; Imperfect-info scaling</p>
          </article>
          <article class="phase-card panel phase-c">
            <div class="phase-badge">C</div>
            <h2>Neural methods</h2>
            <span class="phase-meta">Steps 5–6 &middot; May–early June</span>
            <p>Deep CFR / DREAM &middot; Pluribus &rarr; ReBeL &rarr; Student of Games</p>
          </article>
          <article class="phase-card panel phase-d">
            <div class="phase-badge">D</div>
            <h2>Opponent modelling</h2>
            <span class="phase-meta">Steps 7–8 &middot; June–mid July</span>
            <p>Behavioural inference &middot; Safe exploitation theory and search</p>
          </article>
          <article class="phase-card wide panel phase-e">
            <div class="phase-badge">E</div>
            <h2>Multi-agent dynamics</h2>
            <span class="phase-meta">Steps 9–11 &middot; mid July–August</span>
            <p>MARL &middot; Population-based training &middot; Coalition formation</p>
          </article>
          <article class="phase-card wide panel phase-f">
            <div class="phase-badge">F</div>
            <h2>Data-driven approaches</h2>
            <span class="phase-meta">Steps 12–13 &middot; late August–September</span>
            <p>Sequence models &middot; LLM agents &middot; Behavioural analysis pipelines</p>
          </article>
          <article class="phase-card wide panel phase-g">
            <div class="phase-badge">G</div>
            <h2>Integration</h2>
            <span class="phase-meta">Steps 14–15 &middot; late September–mid October</span>
            <p>Evaluation frameworks &middot; Research frontier mapping</p>
          </article>
        </section>
        <p class="caption phase-footer">All phases conclude before the contribution work begins in earnest.</p>
        """,
        slide_no=11,
        accent="blue",
    )

    s12 = shell(
        "Expected outcomes",
        """
        <section class="columns">
          <article class="outcome panel primary">
            <h2>REALISTIC TARGET</h2>
            <h3>What the thesis will deliver</h3>
            <ul>
              <li>Validated adaptation methods</li>
              <li>Characterised safety heuristics</li>
              <li>Working evaluation framework</li>
              <li>Failure modes documented</li>
            </ul>
          </article>
          <article class="outcome panel light">
            <h2>STRETCH DIRECTION</h2>
            <h3>If the work goes deeper</h3>
            <ul>
              <li>Theoretical extensions of the safety theorem</li>
              <li>Novel combinations across contributions</li>
              <li>Surprising empirical results</li>
            </ul>
          </article>
        </section>
        <p class="caption slide-bottom-note">
          Grounded outcomes already constitute a coherent thesis. No pivot required.
        </p>
        """,
        slide_no=12,
        accent="",
    )

    illustration6 = asset_img(
        "illustration6.png",
        "Transformation graphic: fixed-strategy systems becoming adaptive, safe, accountable agents",
        "title-illustration-image",
    )
    if illustration6:
        s13_visual = f'<section class="closing-illustration panel illustration-image-panel">{illustration6}</section>'
    else:
        s13_visual = """
          <section class="illustration-placeholder closing-illustration panel">
            <div>
              <span class="placeholder-label">illustration6</span>
              <span class="placeholder-meta">Transformation graphic: fixed-strategy systems becoming adaptive, safe, accountable agents</span>
            </div>
          </section>
        """
    s13 = shell(
        "Today to Tomorrow",
        f"""
        <section class="closing-layout">
          {s13_visual}
          <p class="thank-you">Thank you. Questions welcome.</p>
        </section>
        """,
        slide_no=13,
        accent="",
    )

    return [
        ("slide_01_title.html", s1),
        ("slide_02_why_now.html", s2),
        ("slide_03_use_cases.html", s3),
        ("slide_04_games.html", s4),
        ("slide_05_state_of_art.html", s5),
        ("slide_06_limitation.html", s6),
        ("slide_07_open_problems.html", s7),
        ("slide_08_contribution_1.html", s8),
        ("slide_09_contribution_2.html", s9),
        ("slide_10_contribution_3.html", s10),
        ("slide_11_study_path.html", s11),
        ("slide_12_outcomes.html", s12),
        ("slide_13_close.html", s13),
    ]


def write_html_files() -> list[Path]:
    write_common_css()
    HTML_DIR.mkdir(parents=True, exist_ok=True)
    paths = []
    for filename, content in slides():
        path = HTML_DIR / filename
        path.write_text(content, encoding="utf-8")
        paths.append(path)
    return paths


def render_pngs(html_paths: list[Path]) -> list[Path]:
    from playwright.sync_api import sync_playwright

    PNG_DIR.mkdir(parents=True, exist_ok=True)
    png_paths: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": WIDTH, "height": HEIGHT}, device_scale_factor=2)
        for index, html_path in enumerate(html_paths, start=1):
            png_path = PNG_DIR / f"slide_{index:02d}.png"
            page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
            page.screenshot(path=str(png_path), full_page=False)
            png_paths.append(png_path)
        browser.close()
    return png_paths


def verify_pngs(png_paths: list[Path]) -> None:
    expected = (WIDTH * 2, HEIGHT * 2)  # device_scale_factor=2 doubles pixel dimensions
    for path in png_paths:
        with Image.open(path) as image:
            if image.size != expected:
                raise RuntimeError(f"{path} has size {image.size}, expected {expected}")
            stat = ImageStat.Stat(image.convert("L"))
            if stat.stddev[0] < 2.0:
                raise RuntimeError(f"{path} appears blank or nearly blank")


def build_pptx(png_paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for path in png_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)


def parse_speaker_notes() -> dict[int, str]:
    if not SCRIPT_PATH.exists():
        return {}

    text = SCRIPT_PATH.read_text(encoding="utf-8")
    headings = list(re.finditer(r"^## Slide\s+(\d+)\s+[—-]\s+(.+)$", text, flags=re.MULTILINE))
    notes: dict[int, str] = {}

    for index, heading in enumerate(headings):
        slide_no = int(heading.group(1))
        section_end = headings[index + 1].start() if index + 1 < len(headings) else len(text)
        section = text[heading.end() : section_end]
        say_match = re.search(r"^\*\*Say:\*\*\s*$", section, flags=re.MULTILINE)
        if not say_match:
            continue

        title = re.sub(r"\s*\([^)]*\)\s*$", "", heading.group(2)).strip()
        raw_note = section[say_match.end() :].split("\n---", 1)[0].strip()
        note = markdown_to_plain_note(raw_note)
        if note:
            notes[slide_no] = f"{title}\n\n{note}"

    return notes


def markdown_to_plain_note(text: str) -> str:
    lines = []
    for line in text.splitlines():
        clean = line.strip()
        clean = re.sub(r"^[-*]\s+", "", clean)
        clean = re.sub(r"\*\*(.*?)\*\*", r"\1", clean)
        clean = re.sub(r"\*(.*?)\*", r"\1", clean)
        lines.append(clean)
    return "\n".join(lines).strip()


def inject_speaker_notes(pptx_path: Path, notes: dict[int, str]) -> None:
    if not notes:
        return

    flatpak_spawn = shutil.which("flatpak-spawn")
    if flatpak_spawn:
        soffice_cmd = [flatpak_spawn, "--host", "/usr/bin/soffice"]
        python_cmd = [
            flatpak_spawn,
            "--host",
            "env",
            "PYTHONPATH=/usr/lib64/python3.14/site-packages:/usr/lib/python3/dist-packages:"
            "/usr/lib64/libreoffice/program:/usr/lib/libreoffice/program",
            "python3",
        ]
    elif shutil.which("soffice"):
        soffice_cmd = ["soffice"]
        python_cmd = ["python3"]
    elif shutil.which("libreoffice"):
        soffice_cmd = ["libreoffice"]
        python_cmd = ["python3"]
    else:
        raise RuntimeError("LibreOffice is required to add speaker notes to the PPTX.")

    helper_dir = HERE / ".tmp"
    helper_dir.mkdir(parents=True, exist_ok=True)
    notes_json = helper_dir / "speaker_notes.json"
    helper_script = helper_dir / "add_speaker_notes_uno.py"
    notes_json.write_text(json.dumps(notes, ensure_ascii=False), encoding="utf-8")
    helper_script.write_text(
        textwrap.dedent(
            """
            import json
            import sys
            import time
            from pathlib import Path

            import uno
            from com.sun.star.beans import PropertyValue


            def prop(name, value):
                item = PropertyValue()
                item.Name = name
                item.Value = value
                return item


            pptx_path = Path(sys.argv[1]).resolve()
            notes_path = Path(sys.argv[2]).resolve()
            notes = {int(key): value for key, value in json.loads(notes_path.read_text(encoding="utf-8")).items()}

            ctx = uno.getComponentContext()
            resolver = ctx.ServiceManager.createInstanceWithContext(
                "com.sun.star.bridge.UnoUrlResolver", ctx
            )
            remote = None
            last_error = None
            for _ in range(40):
                try:
                    remote = resolver.resolve(
                        "uno:socket,host=localhost,port=21073;urp;StarOffice.ComponentContext"
                    )
                    break
                except Exception as exc:
                    last_error = exc
                    time.sleep(0.25)
            if remote is None:
                raise RuntimeError(f"Could not connect to LibreOffice: {last_error}")

            desktop = remote.ServiceManager.createInstanceWithContext(
                "com.sun.star.frame.Desktop", remote
            )
            doc = desktop.loadComponentFromURL(
                pptx_path.as_uri(),
                "_blank",
                0,
                (prop("Hidden", True), prop("ReadOnly", False)),
            )
            if doc is None:
                raise RuntimeError(f"LibreOffice could not load {pptx_path}")

            try:
                pages = doc.getDrawPages()
                for slide_no, note in notes.items():
                    if not 1 <= slide_no <= pages.getCount():
                        continue

                    notes_page = pages.getByIndex(slide_no - 1).getNotesPage()
                    target = None
                    for index in range(notes_page.getCount()):
                        shape = notes_page.getByIndex(index)
                        if getattr(shape, "ShapeType", "") == "com.sun.star.presentation.NotesShape":
                            target = shape
                            break
                    if target is None:
                        for index in range(notes_page.getCount()):
                            shape = notes_page.getByIndex(index)
                            if hasattr(shape, "String"):
                                target = shape
                                break
                    if target is None:
                        target = doc.createInstance("com.sun.star.drawing.TextShape")
                        target.Position = uno.createUnoStruct("com.sun.star.awt.Point", 914400, 3657600)
                        target.Size = uno.createUnoStruct("com.sun.star.awt.Size", 5486400, 4572000)
                        notes_page.add(target)

                    target.String = note
                doc.store()
            finally:
                doc.close(True)
            """
        ).strip()
        + "\n",
        encoding="utf-8",
    )

    profile = f"file:///tmp/rlplan-lo-profile-{pptx_path.stem}"
    server = subprocess.Popen(
        [
            *soffice_cmd,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "-env:UserInstallation=" + profile,
            "--accept=socket,host=localhost,port=21073;urp;StarOffice.ComponentContext",
            "--nodefault",
            "--norestore",
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    try:
        result = subprocess.run(
            [*python_cmd, str(helper_script.resolve()), str(pptx_path.resolve()), str(notes_json.resolve())],
            check=False,
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(
                "LibreOffice failed to add speaker notes.\n"
                f"stdout:\n{result.stdout}\n"
                f"stderr:\n{result.stderr}"
            )
    finally:
        server.terminate()
        try:
            server.communicate(timeout=5)
        except subprocess.TimeoutExpired:
            server.kill()
            server.communicate()
        notes_json.unlink(missing_ok=True)
        helper_script.unlink(missing_ok=True)
        try:
            helper_dir.rmdir()
        except OSError:
            pass


def build() -> None:
    html_paths = write_html_files()
    png_paths = render_pngs(html_paths)
    verify_pngs(png_paths)
    build_pptx(png_paths)
    speaker_notes = parse_speaker_notes()
    inject_speaker_notes(OUTPUT, speaker_notes)
    print(f"Wrote {len(html_paths)} HTML slides to {HTML_DIR}")
    print(f"Wrote {len(png_paths)} PNG slides to {PNG_DIR}")
    print(f"Added speaker notes to {len(speaker_notes)} slides from {SCRIPT_PATH}")
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
